# -*- coding: utf-8 -*-
import os
import json
import time
import threading
import requests
import datetime
import schedule
import textwrap
from dataclasses import dataclass
from typing import Any, Callable, Dict, List, Literal, Optional, TypedDict, Tuple
from collections.abc import Mapping, Sequence
from langchain_openai import ChatOpenAI
from langchain_core.messages import AIMessage, HumanMessage, SystemMessage
try:
    from langgraph.graph import END, START, StateGraph
    _LANGGRAPH_AVAILABLE = True
except Exception:
    END = None
    START = None
    StateGraph = None
    _LANGGRAPH_AVAILABLE = False

# 新增：本地知识库相关导入
import re
from pathlib import Path
import docx
from docx.document import Document
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.table import _Cell, Table
from docx.text.paragraph import Paragraph
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# 用于处理 .doc 文件的库
try:
    import win32com.client
    WIN32COM_AVAILABLE = True
except ImportError:
    WIN32COM_AVAILABLE = False

from utils import util
import utils.config_util as cfg
from genagents.genagents import GenerativeAgent
from genagents.modules.memory_stream import ConceptNode, generate_importance_score
from simulation_engine.gpt_structure import get_text_embedding
from urllib3.exceptions import InsecureRequestWarning
from scheduler.thread_manager import MyThread
from core import content_db
from core import stream_manager
from core import member_db
from faymcp import runtime_bridge as mcp_runtime
from llm.execution_manager import (
    ExecutionManager, ExecutionState, ExecutionStatus,
    get_execution_manager, _get_llm_instance,
)

# 加载配置
cfg.load_config()

# 禁用不安全请求警告
requests.packages.urllib3.disable_warnings(category=InsecureRequestWarning)

agents = {}  # type: dict[str, GenerativeAgent]
agent_lock = threading.RLock()  # 使用可重入锁保护agent对象
reflection_lock = threading.RLock()  # 使用可重入锁保护reflection_time
save_lock = threading.RLock()  # 使用可重入锁保护save_time
reflection_time = None
save_time = None

memory_cleared = False  # 添加记忆清除标记
# 新增: 当前会话用户名及按用户获取memory目录的辅助函数
current_username = None  # 当前会话用户名

def _log_prompt(messages: List[SystemMessage | HumanMessage | AIMessage], tag: str = ""):
    """No-op placeholder for prompt logging (disabled)."""
    return


def _normalize_short_greeting_text(content: Any) -> str:
    if content is None:
        return ""
    text = content if isinstance(content, str) else str(content)
    text = text.strip().lower()
    if not text:
        return ""
    return re.sub(r"[\s`~!@#$%^&*()\-_=+\[\]{}\\|;:'\",<.>/?，。！？、；：‘’“”（）【】《》…·～]+", "", text)


def _is_current_only_turn(content: Any, observation: Any = None) -> bool:
    normalized = _normalize_short_greeting_text(content)
    if not normalized or len(normalized) > 8:
        return False

    short_greetings = {
        "你好",
        "您好",
        "你好呀",
        "你好啊",
        "嗨",
        "嗨嗨",
        "哈喽",
        "哈啰",
        "哈咯",
        "hello",
        "hi",
        "hey",
        "在吗",
        "在嘛",
        "在么",
        "在不在",
        "忙吗",
        "早",
        "早安",
        "早上好",
        "午安",
        "中午好",
        "下午好",
        "晚上好",
        "晚安",
    }
    return normalized in short_greetings


# 小模型实例（流式，面向用户的快速回复）
llm = _get_llm_instance("small", streaming=True)


@dataclass
class WorkflowToolSpec:
    name: str
    description: str
    schema: Dict[str, Any]
    executor: Callable[[Dict[str, Any], int], Tuple[bool, Optional[str], Optional[str]]]
    example_args: Dict[str, Any]


class ToolCall(TypedDict):
    name: str
    args: Dict[str, Any]


class ToolResult(TypedDict, total=False):
    call: ToolCall
    success: bool
    output: Optional[str]
    error: Optional[str]
    attempt: int


class ConversationMessage(TypedDict):
    role: Literal["user", "assistant"]
    content: str


class AgentState(TypedDict, total=False):
    request: str
    messages: List[ConversationMessage]
    tool_results: List[ToolResult]
    next_action: Optional[ToolCall]
    status: Literal["planning", "needs_tool", "completed", "failed"]
    final_response: Optional[str]
    final_messages: Optional[List[SystemMessage | HumanMessage | AIMessage]]
    _response_streamed: Optional[bool]
    planner_preview: Optional[str]
    audit_log: List[str]
    context: Dict[str, Any]
    error: Optional[str]
    max_steps: int


def _find_last_safe_punct(text: str, punctuation_list) -> int:
    """在文本中查找最后一个安全的标点切分位置，跳过数字中的小数点（如0.85）"""
    last_punct_pos = -1
    for punct in punctuation_list:
        pos = text.rfind(punct)
        # 对英文句点，跳过数字间的小数点
        while pos > 0 and punct == ".":
            prev_ch = text[pos - 1] if pos > 0 else ""
            next_ch = text[pos + 1] if pos + 1 < len(text) else ""
            if prev_ch.isdigit() and next_ch.isdigit():
                pos = text.rfind(punct, 0, pos)
            else:
                break
        if pos > last_punct_pos:
            last_punct_pos = pos
    return last_punct_pos


def _truncate_text(text: Any, limit: int = 400) -> str:
    text_str = "" if text is None else str(text)
    if len(text_str) <= limit:
        return text_str
    return text_str[:limit] + "..."


def _extract_text_from_result(value: Any, *, depth: int = 0) -> List[str]:
    """Try to pull human-readable text snippets from tool results."""
    if value is None:
        return []
    if depth > 10:
        return []

    # 处理字符串 - 尝试解析为 JSON
    if isinstance(value, str):
        text = value.strip()
        if not text:
            return []
        # 尝试解析 JSON 字符串
        if text.startswith('{') or text.startswith('['):
            try:
                parsed = json.loads(text)
                return _extract_text_from_result(parsed, depth=depth + 1)
            except (json.JSONDecodeError, TypeError):
                pass
        return [text]

    if isinstance(value, (int, float, bool)):
        return [str(value)]

    # 处理字典/映射
    if isinstance(value, Mapping):
        # 优先提取 document 字段（知识库查询结果常用）
        if "document" in value:
            doc = value["document"]
            if isinstance(doc, str) and doc.strip():
                return [doc.strip()]

        # 提取 text 字段
        if "text" in value:
            text_val = value["text"]
            if isinstance(text_val, str):
                # 尝试解析嵌套的 JSON
                text_str = text_val.strip()
                if text_str.startswith('{') or text_str.startswith('['):
                    try:
                        parsed = json.loads(text_str)
                        return _extract_text_from_result(parsed, depth=depth + 1)
                    except (json.JSONDecodeError, TypeError):
                        pass
                if text_str:
                    return [text_str]
            else:
                return _extract_text_from_result(text_val, depth=depth + 1)

        # 处理 content 字段（MCP 工具常用格式）
        if "content" in value:
            content = value["content"]
            if isinstance(content, list):
                segments: List[str] = []
                for item in content:
                    segments.extend(_extract_text_from_result(item, depth=depth + 1))
                if segments:
                    return segments
            else:
                return _extract_text_from_result(content, depth=depth + 1)

        # 处理 results 字段（查询结果常用）
        if "results" in value:
            results = value["results"]
            if isinstance(results, list):
                segments: List[str] = []
                for item in results:
                    segments.extend(_extract_text_from_result(item, depth=depth + 1))
                if segments:
                    return segments

        # 遍历其他字段
        segments = []
        skip_keys = {"meta", "annotations", "uid", "id", "messageId", "type", "distance", "metadata", "count", "isError", "structuredContent"}
        for key, item in value.items():
            if key in skip_keys:
                continue
            item_segments = _extract_text_from_result(item, depth=depth + 1)
            for seg in item_segments:
                segments.append(f"{key}: {seg}")
        return segments

    # 处理列表/序列
    if isinstance(value, Sequence) and not isinstance(value, (str, bytes, bytearray)):
        segments: List[str] = []
        for item in value:
            segments.extend(_extract_text_from_result(item, depth=depth + 1))
        return segments

    # 处理有 content 属性的对象（优先处理，因为 MCP 返回的对象通常有 content）
    if hasattr(value, "content") and not callable(getattr(value, "content")):
        content = getattr(value, "content", None)
        if content is not None:
            return _extract_text_from_result(content, depth=depth + 1)

    # 处理有 text 属性的对象（如 TextContent）
    if hasattr(value, "text") and not callable(getattr(value, "text")):
        text = getattr(value, "text", "")
        if isinstance(text, str):
            text_str = text.strip()
            # 尝试解析嵌套的 JSON
            if text_str.startswith('{') or text_str.startswith('['):
                try:
                    parsed = json.loads(text_str)
                    return _extract_text_from_result(parsed, depth=depth + 1)
                except (json.JSONDecodeError, TypeError):
                    pass
            if text_str:
                return [text_str]
        return _extract_text_from_result(text, depth=depth + 1)

    # 处理有 __dict__ 的对象
    if hasattr(value, "__dict__"):
        obj_dict = vars(value)
        # 跳过无用字段
        skip_attrs = {"meta", "annotations", "type", "isError", "structuredContent"}
        filtered_dict = {k: v for k, v in obj_dict.items() if k not in skip_attrs and v is not None}
        if filtered_dict:
            return _extract_text_from_result(filtered_dict, depth=depth + 1)

    # 最后尝试转字符串（但避免输出类似 "TextContent(...)" 的格式）
    text = str(value).strip()
    # 过滤掉看起来像对象表示的字符串
    if text and not text.startswith('<') and '=' not in text[:50]:
        return [text]
    return []


def _normalize_tool_output(result: Any) -> str:
    """Convert structured tool output to a concise human-readable string."""
    if result is None:
        return ""

    segments = _extract_text_from_result(result)
    if segments:
        # 过滤空字符串，去重，拼接
        cleaned = []
        seen = set()
        for segment in segments:
            if segment and segment not in seen:
                seen.add(segment)
                cleaned.append(segment)
        if cleaned:
            return "\n".join(cleaned)

    # 如果提取失败，尝试返回简化的 JSON
    try:
        return json.dumps(result, ensure_ascii=False, default=lambda o: getattr(o, "__dict__", str(o)))
    except TypeError:
        return str(result)


def _apply_question_placeholder(value: Any, question: str) -> Any:
    """Recursively replace question placeholder inside params."""
    if isinstance(value, str):
        # 兼容 {question} 与 {{question}} 两种写法
        return value.replace("{{question}}", question).replace("{question}", question)
    if isinstance(value, Mapping):
        return {k: _apply_question_placeholder(v, question) for k, v in value.items()}
    if isinstance(value, list):
        return [_apply_question_placeholder(item, question) for item in value]
    return value


def _remove_prestart_from_text(text: str, keep_marked: bool = True) -> str:
    """从文本中处理 prestart 标签

    Args:
        text: 输入文本
        keep_marked: 如果为True，保留 keep="true" 的 prestart 完整标签和内容
                    如果为False，移除所有 prestart 内容

    处理逻辑：
    - <prestart>...</prestart> - 完全移除（标签和内容）
    - <prestart keep="true">...</prestart> - 保留完整标签和内容（如果 keep_marked=True）
    """
    if not text:
        return text
    import re

    if keep_marked:
        # 只移除没有 keep="true" 属性的 prestart 标签及其内容
        # 保留 keep="true" 的完整标签
        text = re.sub(r'<prestart>[\s\S]*?</prestart>', '', text, flags=re.IGNORECASE)
    else:
        # 移除所有 prestart 内容（包括有 keep 属性的）
        text = re.sub(r'<prestart[^>]*>[\s\S]*?</prestart>', '', text, flags=re.IGNORECASE)

    return text.strip()


def _remove_think_from_text(text: str) -> str:
    """从文本中移除 think 标签及其内容"""
    if not text:
        return text
    import re
    cleaned = re.sub(r'<think>[\s\S]*?</think>', '', text, flags=re.IGNORECASE)
    cleaned = re.sub(r'</?think>', '', cleaned, flags=re.IGNORECASE)
    return cleaned.strip()


def _strip_json_code_fence(text: str) -> str:
    """Strip ```json ... ``` wrappers if present."""
    if not text:
        return text
    import re
    trimmed = text.strip()
    match = re.match(r"^```(?:json)?\\s*(.*?)\\s*```$", trimmed, flags=re.IGNORECASE | re.DOTALL)
    if match:
        return match.group(1).strip()
    return text


def _format_conversation_block(conversation: List[Dict], username: str = "User") -> str:
    """格式化对话记录，每条消息用代码块包裹

    格式示例：
    主人：
    ```
    你好
    ```
    Fay：
    ```
    主人，你又想起我来了？
    ```

    支持消息中携带 username 字段来区分不同用户（多用户场景）
    """
    if not conversation:
        return "（暂无对话）"

    formatted_lines = []

    for msg in conversation:
        role = msg.get('role', '')
        content = msg.get('content', '')
        msg_username = msg.get('username', '')  # 消息携带的用户名

        # 移除 prestart 标签（保留 keep="true" 的内容）和 think 标签
        content = _remove_prestart_from_text(content)
        content = _remove_think_from_text(content)

        # 跳过空内容的消息
        if not content or not content.strip():
            continue

        # 根据角色显示名称
        if role.lower() in ('user', 'human', '用户'):
            # 优先使用消息中的用户名，其次使用传入的默认用户名
            actual_username = msg_username or username
            role_name = "主人" if actual_username == "User" else actual_username
        elif role.lower() in ('assistant', 'ai', 'fay'):
            role_name = "Fay"
        else:
            role_name = role

        formatted_lines.append(f"{role_name}：\n```\n{content}\n```")

    return "\n".join(formatted_lines) if formatted_lines else "（暂无对话）"


def _run_prestart_tools(user_question: str) -> List[Dict[str, Any]]:
    """Call configured prestart MCP tools and return a list of result objects."""
    try:
        tools = mcp_runtime.list_runnable_prestart_tools()
    except Exception as exc:
        util.log(1, f"获取预启动工具列表失败: {exc}")
        return []
    if not tools:
        return []

    results: List[Dict[str, Any]] = []
    for item in tools:
        server_id = item.get("server_id")
        tool_name = item.get("tool")
        if not server_id or not tool_name:
            continue
        params = item.get("params") or {}
        include_history = item.get("include_history", True)

        try:
            filled_params = _apply_question_placeholder(params, user_question)
        except Exception:
            filled_params = params or {}

        try:
            success, result = mcp_runtime.call_tool(
                int(server_id),
                tool_name,
                filled_params,
                skip_enabled_check=True,
            )
        except Exception as exc:
            util.log(1, f"预启动工具 {tool_name} 调用异常: {exc}")
            continue

        if success:
            output = _normalize_tool_output(result)
            if output and output.strip():
                # 格式化参数显示
                params_str = ""
                if filled_params:
                    try:
                        # 将参数格式化为 (key=value, ...)
                        items = [f"{k}={v}" for k, v in filled_params.items()]
                        params_str = f"({', '.join(items)})"
                    except Exception:
                        pass

                formatted_output = f"【{tool_name}】{params_str}\n{output.strip()}"
                results.append({
                    "text": formatted_output,
                    "include_history": include_history
                })
        else:
            error_msg = str(result) if result is not None else "未知错误"
            util.log(1, f"预启动工具 {tool_name} 执行失败: {error_msg}")

    return results


def _truncate_history(
    history: List[ToolResult],
    limit: Optional[int] = None,
    output_limit: Optional[int] = None,
) -> str:
    if not history:
        return "（暂无）"
    lines: List[str] = []
    selected = history if limit is None else history[-limit:]
    for item in selected:
        call = item.get("call", {})
        name = call.get("name", "未知工具")
        attempt = item.get("attempt", 0)
        success = item.get("success", False)
        status = "成功" if success else "失败"
        lines.append(f"- {name} 第 {attempt} 次 → {status}")
        output = item.get("output")
        if output is not None:
            output_text = str(output)
            if output_limit is not None:
                output_text = _truncate_text(output_text, output_limit)
            lines.append("  输出：" + output_text)
        error = item.get("error")
        if error is not None:
            error_text = str(error)
            if output_limit is not None:
                error_text = _truncate_text(error_text, output_limit)
            lines.append("  错误：" + error_text)
    return "\n".join(lines)


def _format_schema_parameters(schema: Dict[str, Any]) -> List[str]:
    if not schema:
        return ["  - 无参数"]
    props = schema.get("properties") or {}
    if not props:
        return ["  - 无参数"]
    required = set(schema.get("required") or [])
    lines: List[str] = []
    for field, meta in props.items():
        meta = meta or {}
        field_type = meta.get("type", "string")
        desc = (meta.get("description") or "").strip()
        req_label = "必填" if field in required else "可选"
        line = f"  - {field} ({field_type}，{req_label})"
        if desc:
            line += f"：{desc}"
        lines.append(line)
    return lines or ["  - 无参数"]


def _generate_example_args(schema: Dict[str, Any]) -> Dict[str, Any]:
    example: Dict[str, Any] = {}
    if not schema:
        return example
    props = schema.get("properties") or {}
    for field, meta in props.items():
        meta = meta or {}
        if "default" in meta:
            example[field] = meta["default"]
            continue
        enum_values = meta.get("enum") or []
        if enum_values:
            example[field] = enum_values[0]
            continue
        field_type = meta.get("type", "string")
        if field_type in ("number", "integer"):
            example[field] = 0
        elif field_type == "boolean":
            example[field] = True
        elif field_type == "array":
            example[field] = []
        elif field_type == "object":
            example[field] = {}
        else:
            description_hint = meta.get("description") or ""
            example[field] = description_hint or ""
    return example


def _format_tool_block(spec: WorkflowToolSpec) -> str:
    param_lines = _format_schema_parameters(spec.schema)
    example = json.dumps(spec.example_args, ensure_ascii=False) if spec.example_args else "{}"
    lines = [
        f"- 工具名：{spec.name}",
        f"  功能：{spec.description or '暂无描述'}",
        "  参数：",
        *param_lines,
        f"  示例：{example}",
    ]
    return "\n".join(lines)


def _build_workflow_tool_spec(tool_def: Dict[str, Any]) -> Optional[WorkflowToolSpec]:
    if not tool_def:
        return None
    name = tool_def.get("name")
    server_id = tool_def.get("server_id")
    if not name:
        return None
    if server_id is None:
        util.log(1, f"工具 {name} 缺少 server_id，跳过该工具")
        return None
    try:
        server_id = int(server_id)
    except Exception:
        util.log(1, f"工具 {name} 的 server_id 无效: {server_id}")
        return None
    description = tool_def.get("description") or tool_def.get("summary") or ""
    schema = tool_def.get("inputSchema") or {}
    example_args = _generate_example_args(schema)

    def _executor(args: Dict[str, Any], attempt: int) -> Tuple[bool, Optional[str], Optional[str]]:
        try:
            success, result = mcp_runtime.call_tool(server_id, name, args)
        except Exception as exc:
            util.log(1, f"调用工具 {name} 异常: {exc}")
            return False, None, str(exc)

        if success:
            output = _normalize_tool_output(result)
            return True, output, None

        error_msg = str(result) if result is not None else "未知错误"
        util.log(1, f"调用工具 {name} 失败: {error_msg}")
        return False, None, error_msg

    return WorkflowToolSpec(
        name=name,
        description=description,
        schema=schema,
        executor=_executor,
        example_args=example_args,
    )


def _format_tools_for_prompt(tool_specs: Dict[str, WorkflowToolSpec]) -> str:
    if not tool_specs:
        return "（暂无可用工具）"
    return "\n".join(_format_tool_block(spec) for spec in tool_specs.values())


def _merge_system_input(*segments: str) -> str:
    parts = []
    for segment in segments:
        if isinstance(segment, str):
            cleaned = segment.strip()
            if cleaned:
                parts.append(cleaned)
    return "\n\n---\n\n".join(parts)


def _format_context_section(title: str, content: Any, empty_placeholder: Optional[str] = None) -> str:
    cleaned = ""
    if isinstance(content, str):
        cleaned = content.strip()
    elif content is not None:
        cleaned = str(content).strip()
    if not cleaned:
        if empty_placeholder is None:
            return ""
        cleaned = empty_placeholder
    return f"**{title}**\n{cleaned}"


def _get_latest_tool_result_text(history: List[ToolResult]) -> str:
    if not history:
        return ""
    for item in reversed(history):
        if not isinstance(item, dict):
            continue
        call = item.get("call", {}) or {}
        name = call.get("name") or "unknown_tool"
        args = call.get("args") or {}
        args_text = json.dumps(args, ensure_ascii=False) if args else "{}"
        if item.get("success") and item.get("output"):
            return "\n".join(
                [
                    f"工具：{name}",
                    f"参数：{args_text}",
                    f"结果：{item.get('output')}",
                ]
            )
        if item.get("error"):
            return "\n".join(
                [
                    f"工具：{name}",
                    f"参数：{args_text}",
                    f"错误：{item.get('error')}",
                ]
            )
    return ""


def _is_tool_trace_message(content: str) -> bool:
    if not isinstance(content, str):
        return False
    stripped = content.strip()
    if not stripped:
        return False
    return stripped.startswith("[PLAN]") or stripped.startswith("[TOOL]")


def _build_dialogue_messages(
    conversation: List[Dict[str, Any]],
    username: str,
    fallback_request: str = "",
) -> List[HumanMessage | AIMessage]:
    if not conversation:
        fallback_text = fallback_request.strip()
        return [HumanMessage(content=fallback_text)] if fallback_text else []

    user_names = set()
    for msg in conversation:
        role = str(msg.get("role", "")).lower()
        if role not in ("user", "human", "用户"):
            continue
        actual_username = msg.get("username", "") or username
        display_name = "主人" if actual_username == "User" else str(actual_username)
        user_names.add(display_name)
    need_user_label = len(user_names) > 1

    messages: List[HumanMessage | AIMessage] = []
    for msg in conversation:
        role = str(msg.get("role", "")).lower()
        content = msg.get("content", "")
        content = _remove_prestart_from_text(content)
        content = _remove_think_from_text(content)
        if not content or not content.strip():
            continue
        cleaned = content.strip()
        if _is_tool_trace_message(cleaned):
            continue

        if role in ("assistant", "ai", "fay"):
            messages.append(AIMessage(content=cleaned))
            continue

        if role in ("user", "human", "用户"):
            actual_username = msg.get("username", "") or username
            display_name = "主人" if actual_username == "User" else str(actual_username)
            human_content = f"{display_name}：{cleaned}" if need_user_label else cleaned
            messages.append(HumanMessage(content=human_content))
            continue

        messages.append(HumanMessage(content=f"{msg.get('role', 'unknown')}：{cleaned}"))

    if not messages or not any(isinstance(m, HumanMessage) for m in messages):
        fallback_text = fallback_request.strip() or "你好"
        messages.append(HumanMessage(content=fallback_text))

    # 规范化消息列表：确保 user/assistant 交替出现，兼容严格的模型模板
    normalized: List[HumanMessage | AIMessage] = []
    for m in messages:
        if normalized and type(normalized[-1]) is type(m):
            # 合并连续同角色消息
            normalized[-1] = type(m)(content=normalized[-1].content + "\n" + m.content)
        else:
            normalized.append(m)
    # 确保第一条消息是 HumanMessage（模型模板要求）
    if normalized and isinstance(normalized[0], AIMessage):
        normalized.insert(0, HumanMessage(content=fallback_request.strip() or "你好"))
    return normalized


def _build_planner_messages(state: AgentState) -> List[SystemMessage | HumanMessage | AIMessage]:
    context = state.get("context", {}) or {}
    system_prompt = context.get("system_prompt", "")
    request = state.get("request", "")
    tool_specs = context.get("tool_registry", {}) or {}
    planner_preview = state.get("planner_preview")
    conversation = state.get("messages", []) or []
    history = state.get("tool_results", []) or []
    memory_context = context.get("memory_context", "")
    observation = context.get("observation", "")
    prestart_context = context.get("prestart_context", "")
    username = context.get("username", "User")

    history_text = _truncate_history(history)
    latest_tool_result_text = _get_latest_tool_result_text(history)
    tools_text = _format_tools_for_prompt(tool_specs)
    if prestart_context and prestart_context.strip():
        formatted_items = []
        for item in prestart_context.split("\n\n"):
            item = item.strip()
            if not item:
                continue
            lines = item.split("\n", 1)
            if len(lines) == 2:
                header, result = lines
                formatted_items.append(f"{header}\n```\n{result.strip()}\n```")
            else:
                formatted_items.append(f"```\n{item}\n```")
        wrapped_results = "\n".join(formatted_items)
        prestart_section = f"\n**预启动工具结果**\n{wrapped_results}\n---\n"
    else:
        prestart_section = ""

    knowledge_hint = context.get("knowledge_hint", "")
    # 构建工具名列表（简短，只列名称）
    tool_names = ", ".join(tool_specs.keys()) if tool_specs else "无"

    planner_system = _merge_system_input(
        "你是一个闲聊判断器。判断用户的话是不是闲聊，请严格输出合法 JSON，不要输出其他内容。",
        '输出格式只有两种：'
        '\n1. 是闲聊: {"action": "finish", "message": "你的回复内容"}'
        '\n2. 不是闲聊: {"action": "tool", "keyword": "提取的搜索关键词"}'
        '\n\n什么是闲聊（输出 finish）：'
        '\n- 打招呼：你好、hi、早上好'
        '\n- 情绪表达：我好开心、今天好累'
        '\n- 感谢道别：谢谢、再见、拜拜'
        '\n- 简单确认：好的、收到、明白了'
        '\n- 对你上一句话的回应：哈哈、对的、没错、说得好'
        '\n\n什么不是闲聊（输出 tool）：'
        '\n- 问任何具体事物/概念：XX是什么、你知道XX吗'
        '\n- 要求查询/获取/阅读内容'
        '\n- 提到任何产品名、项目名、专有名词'
        '\n- 任何你需要查资料才能准确回答的问题'
        '\n- 用户的问题涉及下方"可用工具"或"知识库主题"中的任何内容'
        '\n\nkeyword 提取规则：'
        '\n- keyword 必须是具体的搜索主题词，不能是"再查一下""详细说说"等动作描述'
        '\n- 如果用户消息是指代性的（如"再查一下""继续""详细说说"），从对话历史中找到实际话题作为 keyword'
        '\n\n不确定时 → 输出 tool',
        system_prompt,
        _format_context_section("关联记忆", memory_context),
        _format_context_section("可用工具", tool_names),
        _format_context_section("知识库主题（提到这些主题必须输出 tool）", knowledge_hint),
        _format_context_section("最新工具结果", latest_tool_result_text),
        _format_context_section("预启动工具结果", prestart_context),
        _format_context_section("其他观察", observation),
        _format_context_section("历史工具执行", history_text if history_text != "（暂无）" else ""),
        _format_context_section("规划器预览", planner_preview),
    )
    dialogue_messages = _build_dialogue_messages(conversation, username, fallback_request=request)
    if not dialogue_messages:
        dialogue_messages = [HumanMessage(content=request or "你好")]

    return [SystemMessage(content=planner_system), *dialogue_messages]

def _build_final_messages(state: AgentState) -> List[SystemMessage | HumanMessage | AIMessage]:
    context = state.get("context", {}) or {}
    system_prompt = context.get("system_prompt", "")
    request = state.get("request", "")
    memory_context = context.get("memory_context", "")
    observation = context.get("observation", "")
    prestart_context = context.get("prestart_context", "")
    conversation = state.get("messages", []) or []
    planner_preview = state.get("planner_preview")
    username = context.get("username", "User")

    history_text = _truncate_history(state.get("tool_results", []))
    latest_tool_result_text = _get_latest_tool_result_text(state.get("tool_results", []) or [])
    if prestart_context and prestart_context.strip():
        formatted_items = []
        for item in prestart_context.split("\n\n"):
            item = item.strip()
            if not item:
                continue
            lines = item.split("\n", 1)
            if len(lines) == 2:
                header, result = lines
                formatted_items.append(f"{header}\n```\n{result.strip()}\n```")
            else:
                formatted_items.append(f"```\n{item}\n```")
        wrapped_results = "\n".join(formatted_items)
        prestart_section = f"\n**预启动工具结果**\n{wrapped_results}\n---\n"
    else:
        prestart_section = ""

    tool_grounding_instruction = ""
    if latest_tool_result_text and latest_tool_result_text.strip():
        tool_grounding_instruction = (
            "如果已经有成功的工具结果，必须优先依据最新工具结果直接回答用户问题。"
            "不要忽略工具结果，不要回退成泛化寒暄；如果工具结果不足以完整回答，"
            "要明确说明已知结果与缺失信息。"
        )

    final_system = _merge_system_input(
        system_prompt,
        tool_grounding_instruction,
        _format_context_section("关联记忆", memory_context),
        _format_context_section("最新工具结果", latest_tool_result_text),
        _format_context_section("预启动工具结果", prestart_context),
        _format_context_section("其他观察", observation),
        _format_context_section("工具执行摘要", history_text if history_text != "（暂无）" else ""),
        _format_context_section("规划器建议", planner_preview),
    )
    dialogue_messages = _build_dialogue_messages(conversation, username, fallback_request=request)
    if not dialogue_messages:
        dialogue_messages = [HumanMessage(content=request or "你好")]

    return [SystemMessage(content=final_system), *dialogue_messages]

def _call_planner_llm(
    state: AgentState,
    stream_callback: Optional[Callable[[str], None]] = None,
    on_tool_detected: Optional[Callable[[], None]] = None,
) -> Dict[str, Any]:
    """
    调用规划器 LLM，支持流式输出 finish+message 模式。

    Args:
        state: 当前工作流状态
        stream_callback: 可选的流式回调函数，用于实时输出 message 内容
        on_tool_detected: 可选回调，流式中检测到 tool action 时立即调用（用于提前推送过渡语）

    Returns:
        解析后的决策字典
    """
    messages = _build_planner_messages(state)
    _log_prompt(messages, tag="planner")

    # 如果有流式回调，使用流式模式检测 finish+message
    if stream_callback is not None:
        accumulated = ""
        in_message_mode = False
        in_plaintext_mode = False
        tool_early_notified = False
        message_buffer = ""
        message_closed = False
        escape_next = False

        def _stream_message_text(text: str) -> None:
            nonlocal message_buffer, message_closed, escape_next
            if not text or message_closed:
                return
            out_parts = []
            for ch in text:
                if escape_next:
                    escape_next = False
                    message_buffer += ch
                    out_parts.append(ch)
                    continue
                if ch == "\\":
                    escape_next = True
                    message_buffer += ch
                    out_parts.append(ch)
                    continue
                if ch == '"':
                    message_closed = True
                    break
                message_buffer += ch
                out_parts.append(ch)
            if out_parts:
                stream_callback("".join(out_parts))

        for chunk in llm.stream(messages):
            chunk_text = ""
            if isinstance(chunk, str):
                chunk_text = chunk
            elif isinstance(chunk, dict):
                chunk_text = chunk.get("content", "")
            else:
                chunk_text = getattr(chunk, "content", "") or ""

            if not chunk_text:
                continue

            accumulated += chunk_text

            if not in_message_mode and not in_plaintext_mode:
                # 移除可能的 think 标签前缀
                check_text = _remove_think_from_text(accumulated.strip())
                # 压缩空白，兼容 LLM 返回带换行的 JSON 如 {\n  "action": "finish",...}
                compact = re.sub(r'\s+', '', check_text[:60])

                # 检测是否已经进入 finish+message 模式（用压缩后文本匹配）
                finish_compact_prefix = '{"action":"finish","message":"'
                if compact.startswith(finish_compact_prefix):
                    in_message_mode = True
                    # 从原始 check_text 中定位 message 值的起始位置
                    # 找到 "message" 键后第一个引号内的内容
                    msg_key_pos = check_text.find('"message"')
                    if msg_key_pos >= 0:
                        # 跳过 "message" 后的 : 和空白，找到值的起始引号
                        rest = check_text[msg_key_pos + len('"message"'):]
                        colon_pos = rest.find(':')
                        if colon_pos >= 0:
                            after_colon = rest[colon_pos + 1:].lstrip()
                            if after_colon.startswith('"'):
                                message_start = after_colon[1:]  # 跳过起始引号
                                if message_start:
                                    _stream_message_text(message_start)

                # 检测到 tool 模式 → 通过独立回调推送过渡语，减少用户等待
                tool_compact_prefix = '{"action":"tool"'
                if compact.startswith(tool_compact_prefix) and not in_message_mode and not tool_early_notified:
                    tool_early_notified = True
                    if on_tool_detected:
                        on_tool_detected()
                    # 不进入任何流式模式，后续 chunk 静默累积等完整 JSON

                # 如果不是 finish+message，根据已知 JSON 前缀判断是否为纯文本
                if not in_message_mode:
                    if len(check_text) > 3 and not check_text.startswith('{'):
                        # 不以 { 开头，明确是纯文本
                        in_plaintext_mode = True
                        stream_callback(check_text)
                    elif check_text.startswith('{') and len(compact) > 15:
                        # 以 { 开头，压缩空白后检查是否匹配已知 JSON action 前缀
                        if not compact.startswith('{"action"'):
                            in_plaintext_mode = True
                            stream_callback(check_text)
            elif in_message_mode:
                # 已经在 message 模式，直接流式输出新增内容
                _stream_message_text(chunk_text)
            elif in_plaintext_mode:
                # 纯文本模式，直接流式输出每个 chunk
                stream_callback(chunk_text)

        # 处理完整响应
        trimmed = _remove_think_from_text(accumulated.strip())
        trimmed = _strip_json_code_fence(trimmed)

        if in_message_mode:
            # 提取完整 message 内容，去掉结尾的 "}
            # message_buffer 包含从 "message": " 之后的所有内容
            # 需要去掉结尾的 "}（可能有空格）
            full_message = message_buffer.rstrip()
            # 移除结尾的 "} 或 " }
            if full_message.endswith('"}'):
                full_message = full_message[:-2]
            elif full_message.endswith('" }'):
                full_message = full_message[:-3]
            elif full_message.endswith('"'):
                full_message = full_message[:-1]

            # 处理可能的转义字符
            try:
                # 尝试用 JSON 解码字符串（处理 \n, \" 等转义）
                full_message = json.loads(f'"{full_message}"')
            except (json.JSONDecodeError, TypeError):
                pass

            return {
                "action": "finish",
                "message": full_message,
                "_raw": trimmed,
                "_streamed": True
            }

        # 纯文本模式：LLM 直接返回了文本内容（非 JSON），作为 finish 直接输出
        if in_plaintext_mode:
            return {
                "action": "finish",
                "message": trimmed,
                "_raw": trimmed,
                "_streamed": True
            }

        # 非 finish+message 模式，按原逻辑解析（如 tool 调用等 JSON 响应）
        try:
            decision = json.loads(trimmed)
        except json.JSONDecodeError:
            # JSON 解析失败，将整段文本作为直接回复内容
            return {
                "action": "finish",
                "message": trimmed,
                "_raw": trimmed,
                "_streamed": False
            }
        decision.setdefault("_raw", trimmed)
        if tool_early_notified:
            decision["_tool_early_streamed"] = True
        return decision

    # 无流式回调，使用原有的非流式模式
    response = llm.invoke(messages)
    content = getattr(response, "content", None)
    if not isinstance(content, str):
        raise RuntimeError("规划器返回内容异常，未获得字符串。")
    # 先移除 think 标签，兼容带思考标签的模型（如 DeepSeek R1）
    trimmed = _remove_think_from_text(content.strip())
    trimmed = _strip_json_code_fence(trimmed)
    try:
        decision = json.loads(trimmed)
    except json.JSONDecodeError:
        # JSON 解析失败，将整段文本作为直接回复内容（LLM 直接返回了文本而非 JSON）
        return {
            "action": "finish",
            "message": trimmed,
            "_raw": trimmed,
        }
    decision.setdefault("_raw", trimmed)
    return decision


def _plan_next_action(state: AgentState) -> AgentState:
    context = state.get("context", {}) or {}
    audit_log = list(state.get("audit_log", []))
    history = state.get("tool_results", []) or []
    max_steps = state.get("max_steps", 12)
    if len(history) >= max_steps:
        audit_log.append("规划器：超过最大步数，终止流程。")
        return {
            "status": "failed",
            "audit_log": audit_log,
            "error": "工具调用步数超限",
            "context": context,
        }

    # 从 context 获取流式回调（如果有的话）
    stream_callback = context.get("planner_stream_callback")
    decision = _call_planner_llm(state, stream_callback=stream_callback)
    audit_log.append(f"规划器：决策 -> {decision.get('_raw', decision)}")

    action = decision.get("action")
    if action == "tool":
        tool_name = decision.get("tool")
        tool_registry: Dict[str, WorkflowToolSpec] = context.get("tool_registry", {})
        if tool_name not in tool_registry:
            audit_log.append(f"规划器：未知工具 {tool_name}")
            return {
                "status": "failed",
                "audit_log": audit_log,
                "error": f"未知工具 {tool_name}",
                "context": context,
            }
        args = decision.get("args") or {}

        if history:
            last_entry = history[-1]
            last_call = last_entry.get("call", {}) or {}
            if (
                last_entry.get("success")
                and last_call.get("name") == tool_name
                and (last_call.get("args") or {}) == args
                and last_entry.get("output")
            ):
                recent_attempts = sum(
                    1
                    for item in reversed(history)
                    if item.get("call", {}).get("name") == tool_name
                )
                if recent_attempts >= 1:
                    audit_log.append(
                        "规划器：检测到工具重复调用，使用最新结果产出最终回复。"
                    )
                    final_messages = _build_final_messages(state)
                    _log_prompt(final_messages, tag="final")
                    preview = last_entry.get("output")
                    return {
                        "status": "completed",
                        "planner_preview": preview,
                        "final_response": None,
                        "final_messages": final_messages,
                        "audit_log": audit_log,
                        "context": context,
                    }
        return {
            "next_action": {"name": tool_name, "args": args},
            "status": "needs_tool",
            "audit_log": audit_log,
            "context": context,
        }

    if action in {"finish", "finish_text"}:
        preview = decision.get("message")
        was_streamed = decision.get("_streamed", False)
        used_tool = bool(history)
        # 如果 finish 带有 message，直接使用，不需要再调用最终 LLM
        if action == "finish" and preview:
            audit_log.append(f"规划器：任务完成，直接输出回复: {preview[:50]}...")
            return {
                "status": "completed",
                "planner_preview": preview,
                "final_response": preview,
                "final_messages": None,  # 不需要再调用 LLM
                "audit_log": audit_log,
                "context": context,
                "_response_streamed": was_streamed,  # 标记是否已流式输出
            }
        # finish_text / finish 不带 message，调用最终 LLM 生成回复
        final_messages = _build_final_messages(state)
        audit_log.append("规划器：任务完成，准备调用 LLM 生成最终回复。")
        return {
            "status": "completed",
            "planner_preview": preview,
            "final_response": None,
            "final_messages": final_messages,
            "audit_log": audit_log,
            "context": context,
        }

    raise RuntimeError(f"未知的规划器决策: {decision}")


def _execute_tool(state: AgentState) -> AgentState:
    context = dict(state.get("context", {}) or {})
    action = state.get("next_action")
    if not action:
        return {
            "status": "failed",
            "error": "缺少要执行的工具指令",
            "context": context,
        }

    history = list(state.get("tool_results", []) or [])
    audit_log = list(state.get("audit_log", []) or [])
    conversation = list(state.get("messages", []) or [])

    name = action.get("name")
    args = action.get("args", {})
    tool_registry: Dict[str, WorkflowToolSpec] = context.get("tool_registry", {})
    spec = tool_registry.get(name)
    if not spec:
        return {
            "status": "failed",
            "error": f"未知工具 {name}",
            "context": context,
        }

    attempts = sum(1 for item in history if item.get("call", {}).get("name") == name)
    success, output, error = spec.executor(args, attempts)
    result: ToolResult = {
        "call": {"name": name, "args": args},
        "success": success,
        "output": output,
        "error": error,
        "attempt": attempts + 1,
    }
    history.append(result)
    audit_log.append(f"执行器：{name} 第 {result['attempt']} 次 -> {'成功' if success else '失败'}")

    message_lines = [
        f"[TOOL] {name} {'成功' if success else '失败'}。",
    ]
    if output:
        message_lines.append(f"[TOOL] 输出：{_truncate_text(output, 200)}")
    if error:
        message_lines.append(f"[TOOL] 错误：{_truncate_text(error, 200)}")
    conversation.append({"role": "assistant", "content": "\n".join(message_lines)})

    return {
        "tool_results": history,
        "messages": conversation,
        "next_action": None,
        "audit_log": audit_log,
        "status": "planning",
        "error": error if not success else None,
        "context": context,
    }


def _route_decision(state: AgentState) -> str:
    return "call_tool" if state.get("status") == "needs_tool" else "end"


def _build_workflow_app() -> StateGraph:
    if not _LANGGRAPH_AVAILABLE:
        return None
    graph = StateGraph(AgentState)
    graph.add_node("plan_next", _plan_next_action)
    graph.add_node("call_tool", _execute_tool)
    graph.add_edge(START, "plan_next")
    graph.add_conditional_edges(
        "plan_next",
        _route_decision,
        {
            "call_tool": "call_tool",
            "end": END,
        },
    )
    graph.add_edge("call_tool", "plan_next")
    return graph.compile()


_WORKFLOW_APP = _build_workflow_app() if _LANGGRAPH_AVAILABLE else None

def get_user_memory_dir(username=None):
    """根据配置决定是否按用户名隔离记忆目录"""
    if username is None:
        username = current_username
    base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    mem_base = os.path.join(base_dir, "memory")
    try:
        cfg.load_config()
        isolate = cfg.config["memory"]["isolate_by_user"]
    except Exception:
        isolate = False
    if isolate and username:
        return os.path.join(mem_base, str(username))
    return mem_base

def get_current_time_step(username=None):
    """
    获取当前时间作为time_step
    
    返回:
        int: 当前时间步，从0开始，非真实时间
    """
    global agents
    try:
        # 按用户名选择对应agent，若未指定则退回全局agent
        ag = agents.get(username) if username else None
        if ag and ag.memory_stream and ag.memory_stream.seq_nodes:
            # 如果有记忆节点，则使用最后一个节点的created属性加1
            return int(ag.memory_stream.seq_nodes[-1].created) + 1
        else:
            # 如果没有记忆节点或agent未初始化，则使用0
            return 0
    except Exception as e:
        util.log(1, f"获取time_step时出错: {str(e)}，使用0代替")
        return 0

# 新增：本地知识库相关函数
def read_doc_file(file_path):
    """
    读取doc文件内容
    
    参数:
        file_path: doc文件路径
        
    返回:
        str: 文档内容
    """
    try:
        # 方法1: 使用 win32com.client（Windows系统，推荐用于.doc文件）
        if WIN32COM_AVAILABLE:
            word = None
            doc = None
            try:
                import pythoncom
                pythoncom.CoInitialize()  # 初始化COM组件
                
                word = win32com.client.Dispatch("Word.Application")
                word.Visible = False
                doc = word.Documents.Open(file_path)
                content = doc.Content.Text
                
                # 先保存内容，再尝试关闭
                if content and content.strip():
                    try:
                        doc.Close()
                        word.Quit()
                    except Exception as close_e:
                        util.log(1, f"关闭Word应用程序时出错: {str(close_e)}，但内容已成功提取")
                    
                    try:
                        pythoncom.CoUninitialize()  # 清理COM组件
                    except:
                        pass
                    
                    return content.strip()
                
            except Exception as e:
                util.log(1, f"使用 win32com 读取 .doc 文件失败: {str(e)}")
            finally:
                # 确保资源被释放
                try:
                    if doc:
                        doc.Close()
                except:
                    pass
                try:
                    if word:
                        word.Quit()
                except:
                    pass
                try:
                    pythoncom.CoUninitialize()
                except:
                    pass
        
        # 方法2: 简单的二进制文本提取（备选方案）
        try:
            with open(file_path, 'rb') as f:
                raw_data = f.read()
                # 尝试提取可打印的文本
                text_parts = []
                current_text = ""
                
                for byte in raw_data:
                    char = chr(byte) if 32 <= byte <= 126 or byte in [9, 10, 13] else None
                    if char:
                        current_text += char
                    else:
                        if len(current_text) > 3:  # 只保留长度大于3的文本片段
                            text_parts.append(current_text.strip())
                        current_text = ""
                
                if len(current_text) > 3:
                    text_parts.append(current_text.strip())
                
                # 过滤和清理文本
                filtered_parts = []
                for part in text_parts:
                    # 移除过多的重复字符和无意义的片段
                    if (len(part) > 5 and 
                        not part.startswith('Microsoft') and 
                        not all(c in '0123456789-_.' for c in part) and
                        len(set(part)) > 3):  # 字符种类要多样
                        filtered_parts.append(part)
                
                if filtered_parts:
                    return '\n'.join(filtered_parts)
                    
        except Exception as e:
            util.log(1, f"使用二进制方法读取 .doc 文件失败: {str(e)}")
        
        util.log(1, f"无法读取 .doc 文件 {file_path}，建议转换为 .docx 格式")
        return ""
        
    except Exception as e:
        util.log(1, f"读取doc文件 {file_path} 时出错: {str(e)}")
        return ""

def read_docx_file(file_path):
    """
    读取docx文件内容
    
    参数:
        file_path: docx文件路径
        
    返回:
        str: 文档内容
    """
    try:
        doc = docx.Document(file_path)
        content = []
        
        for element in doc.element.body:
            if isinstance(element, CT_P):
                paragraph = Paragraph(element, doc)
                if paragraph.text.strip():
                    content.append(paragraph.text.strip())
            elif isinstance(element, CT_Tbl):
                table = Table(element, doc)
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        content.append(" | ".join(row_text))
        
        return "\n".join(content)
    except Exception as e:
        util.log(1, f"读取docx文件 {file_path} 时出错: {str(e)}")
        return ""
    
def read_pptx_file(file_path):
    """
    读取pptx文件内容
    
    参数:
        file_path: pptx文件路径
        
    返回:
        str: 演示文稿内容
    """
    if not PPTX_AVAILABLE:
        util.log(1, "python-pptx 库未安装，无法读取 PowerPoint 文件")
        return ""
        
    try:
        prs = Presentation(file_path)
        content = []
        
        for i, slide in enumerate(prs.slides):
            slide_content = [f"第{i+1}页："]
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text.strip())
                    
            if len(slide_content) > 1:  # 有内容才添加
                content.append("\n".join(slide_content))
        
        return "\n\n".join(content)
    except Exception as e:
        util.log(1, f"读取pptx文件 {file_path} 时出错: {str(e)}")
        return ""

def load_local_knowledge_base():
    """
    加载本地知识库内容
    
    返回:
        dict: 文件名到内容的映射
    """
    knowledge_base = {}
    
    # 获取llm/data目录路径
    current_dir = os.path.dirname(os.path.abspath(__file__))
    data_dir = os.path.join(current_dir, "data")
    
    if not os.path.exists(data_dir):
        util.log(1, f"知识库目录不存在: {data_dir}")
        return knowledge_base
    
    # 遍历data目录中的文件
    for file_path in Path(data_dir).iterdir():
        if not file_path.is_file():
            continue
            
        file_name = file_path.name
        file_extension = file_path.suffix.lower()
        
        try:
            if file_extension == '.docx':
                content = read_docx_file(str(file_path))
            elif file_extension == '.doc':
                content = read_doc_file(str(file_path))
            elif file_extension == '.pptx':
                content = read_pptx_file(str(file_path))
            else:
                # 尝试作为文本文件读取
                try:
                    with open(file_path, 'r', encoding='utf-8') as f:
                        content = f.read()
                except UnicodeDecodeError:
                    try:
                        with open(file_path, 'r', encoding='gbk') as f:
                            content = f.read()
                    except UnicodeDecodeError:
                        util.log(1, f"无法解码文件: {file_name}")
                        continue
            
            if content.strip():
                knowledge_base[file_name] = content
                util.log(1, f"成功加载知识库文件: {file_name} ({len(content)} 字符)")
            
        except Exception as e:
            util.log(1, f"加载知识库文件 {file_name} 时出错: {str(e)}")
    
    return knowledge_base

def search_knowledge_base(query, knowledge_base, max_results=3):
    """
    在知识库中搜索相关内容
    
    参数:
        query: 查询内容
        knowledge_base: 知识库字典
        max_results: 最大返回结果数
        
    返回:
        list: 相关内容列表
    """
    if not knowledge_base:
        return []
    
    results = []
    query_lower = query.lower()
    
    # 搜索关键词
    query_keywords = re.findall(r'\w+', query_lower)
    
    for file_name, content in knowledge_base.items():
        content_lower = content.lower()
        
        # 计算匹配度
        score = 0
        matched_sentences = []
        
        # 按句子分割内容
        sentences = re.split(r'[。！？\n]', content)
        
        for sentence in sentences:
            if not sentence.strip():
                continue
                
            sentence_lower = sentence.lower()
            sentence_score = 0
            
            # 计算关键词匹配度
            for keyword in query_keywords:
                if keyword in sentence_lower:
                    sentence_score += 1
            
            # 如果句子有匹配，记录
            if sentence_score > 0:
                matched_sentences.append((sentence.strip(), sentence_score))
                score += sentence_score
        
        # 如果有匹配的内容
        if score > 0:
            # 按匹配度排序句子
            matched_sentences.sort(key=lambda x: x[1], reverse=True)
            
            # 取前几个最相关的句子
            relevant_sentences = [sent[0] for sent in matched_sentences[:5] if sent[0]]
            
            if relevant_sentences:
                results.append({
                    'file_name': file_name,
                    'score': score,
                    'content': '\n'.join(relevant_sentences)
                })
    
    # 按匹配度排序
    results.sort(key=lambda x: x['score'], reverse=True)
    
    return results[:max_results]

# 全局知识库缓存
_knowledge_base_cache = None
_knowledge_base_load_time = None
_knowledge_base_file_times = {}  # 存储文件的最后修改时间

def check_knowledge_base_changes():
    """
    检查知识库文件是否有变化
    
    返回:
        bool: 如果有文件变化返回True，否则返回False
    """
    global _knowledge_base_file_times
    
    # 获取llm/data目录路径
    current_dir = os.path.dirname(os.path.abspath(__file__))
    data_dir = os.path.join(current_dir, "data")
    
    if not os.path.exists(data_dir):
        return False
    
    current_file_times = {}
    
    # 遍历data目录中的文件
    for file_path in Path(data_dir).iterdir():
        if not file_path.is_file():
            continue
        
        file_name = file_path.name
        file_extension = file_path.suffix.lower()
        
        # 只检查支持的文件格式
        if file_extension in ['.docx', '.doc', '.pptx', '.txt'] or file_extension == '':
            try:
                mtime = os.path.getmtime(str(file_path))
                current_file_times[file_name] = mtime
            except OSError:
                continue
    
    # 检查是否有变化
    if not _knowledge_base_file_times:
        # 第一次检查，保存文件时间
        _knowledge_base_file_times = current_file_times
        return True
    
    # 比较文件时间
    if set(current_file_times.keys()) != set(_knowledge_base_file_times.keys()):
        # 文件数量发生变化
        _knowledge_base_file_times = current_file_times
        return True
    
    for file_name, mtime in current_file_times.items():
        if file_name not in _knowledge_base_file_times or _knowledge_base_file_times[file_name] != mtime:
            # 文件被修改
            _knowledge_base_file_times = current_file_times
            return True
    
    return False

def init_knowledge_base():
    """
    初始化知识库，在系统启动时调用
    """
    global _knowledge_base_cache, _knowledge_base_load_time
    
    util.log(1, "初始化本地知识库...")
    _knowledge_base_cache = load_local_knowledge_base()
    _knowledge_base_load_time = time.time()
    
    # 初始化文件修改时间跟踪
    check_knowledge_base_changes()
    
    util.log(1, f"知识库初始化完成，共 {len(_knowledge_base_cache)} 个文件")

def get_knowledge_base():
    """
    获取知识库，使用缓存机制
    
    返回:
        dict: 知识库内容
    """
    global _knowledge_base_cache, _knowledge_base_load_time
    
    # 如果缓存为空，先初始化
    if _knowledge_base_cache is None:
        init_knowledge_base()
        return _knowledge_base_cache
    
    # 检查文件是否有变化
    if check_knowledge_base_changes():
        util.log(1, "检测到知识库文件变化，正在重新加载...")
        _knowledge_base_cache = load_local_knowledge_base()
        _knowledge_base_load_time = time.time()
        util.log(1, f"知识库重新加载完成，共 {len(_knowledge_base_cache)} 个文件")
    
    return _knowledge_base_cache


# 定时保存记忆的线程
def memory_scheduler_thread():
    """
    定时任务线程，运行schedule调度器
    """
    while True:
        schedule.run_pending()
        time.sleep(60)  # 每分钟检查一次是否有定时任务需要执行

# 初始化定时保存记忆的任务
def init_memory_scheduler():
    """
    初始化定时保存记忆的任务
    """
    global agents
    
    # 确保agent已经创建
    agent = None
    if not agents:
        util.log(1, '创建代理实例...')
        agent = create_agent()
    else:
        agent = agents.get("User")
        if agent is None and len(agents) > 0:
            agent = next(iter(agents.values()))

    # 启动阶段做一次 embedding 维度检查，避免首条消息时触发
    try:
        if agent and agent.memory_stream and hasattr(agent.memory_stream, "precheck_embedding_dimensions"):
            result = agent.memory_stream.precheck_embedding_dimensions()
            if result.get("checked"):
                util.log(
                    1,
                    f"启动阶段记忆 embedding 维度检查完成: dim={result.get('expected_dim')}, 修复={result.get('fixed')}"
                )
            else:
                util.log(1, "启动阶段记忆 embedding 维度检查跳过（无记忆/无embedding）")
    except Exception as e:
        util.log(1, f"启动阶段 embedding 维度检查失败: {str(e)}")
    
    # 设置每天0点保存记忆
    schedule.every().day.at("00:00").do(save_agent_memory)

    # 设置每天晚上11点执行反思
    schedule.every().day.at("23:00").do(perform_daily_reflection)

    # 设置执行用户画像分析（测试用11:30，正式改回22:35）
    schedule.every().day.at("11:30").do(perform_user_portrait_analysis)

    # 启动定时任务线程
    scheduler_thread = MyThread(target=memory_scheduler_thread)
    scheduler_thread.start()

    util.log(1, '定时任务已启动：每天0点保存记忆，每天11:30用户画像分析，每天23点执行反思')

def check_memory_files(username=None):
    """
    检查memory目录及其必要文件是否存在
    
    返回:
        memory_dir: memory目录路径
        is_complete: 是否已经存在完整的memory目录结构
    """
    
    # 根据配置与用户名获取memory目录路径
    memory_dir = get_user_memory_dir(username)

    # 检查memory目录是否存在，不存在则创建
    if not os.path.exists(memory_dir):
        os.makedirs(memory_dir)
        util.log(1, f"创建memory目录: {memory_dir}")
    
    # 删除.memory_cleared标记文件（如果存在）
    base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
    mem_base = os.path.join(base_dir, "memory")
    memory_cleared_flag_file = os.path.join(mem_base, ".memory_cleared")
    if os.path.exists(memory_cleared_flag_file):
        try:
            os.remove(memory_cleared_flag_file)
            util.log(1, f"清除删除记忆标记文件: {memory_cleared_flag_file}")
            # 重置记忆清除标记
            global memory_cleared
            memory_cleared = False
        except Exception as e:
            util.log(1, f"清除删除记忆标记文件时出错: {str(e)}")
    
    # 检查meta.json是否存在
    meta_file = os.path.join(memory_dir, "meta.json")
    is_complete = os.path.exists(meta_file)
    
    # 检查memory_stream目录是否存在，不存在则创建
    memory_stream_dir = os.path.join(memory_dir, "memory_stream")
    if not os.path.exists(memory_stream_dir):
        os.makedirs(memory_stream_dir)
        util.log(1, f"创建memory_stream目录: {memory_stream_dir}")
    
    # 检查必要的文件是否存在
    embeddings_path = os.path.join(memory_stream_dir, "embeddings.json")
    nodes_path = os.path.join(memory_stream_dir, "nodes.json")
    
    # 检查文件是否存在且不为空
    is_complete = (os.path.exists(embeddings_path) and os.path.getsize(embeddings_path) > 2 and
                  os.path.exists(nodes_path) and os.path.getsize(nodes_path) > 2)
    
    # 如果文件不存在，创建空的JSON文件
    if not os.path.exists(embeddings_path):
        with open(embeddings_path, 'w', encoding='utf-8') as f:
            f.write('{}')
    
    if not os.path.exists(nodes_path):
        with open(nodes_path, 'w', encoding='utf-8') as f:
            f.write('[]')
    
    return memory_dir, is_complete

def create_agent(username=None):
    """
    创建一个GenerativeAgent实例
    
    返回:
        agent: GenerativeAgent对象
    """
    global agents
    
    if username is None:
        username = "User"
    
    # 创建/复用代理
    with agent_lock:
        if username in agents:
            return agents[username]
        
        memory_dir, is_exist = check_memory_files(username)
        agent = GenerativeAgent(memory_dir)
        
        # 检查是否有scratch属性，如果没有则添加
        if not hasattr(agent, 'scratch'):
            agent.scratch = {}
        
        # 初始化代理的scratch数据，始终从config_util实时加载
        scratch_data = {
            "first_name": cfg.config["attribute"]["name"],
            "last_name": "",
            "age": cfg.config["attribute"]["age"],
            "sex": cfg.config["attribute"]["gender"],
            "additional": cfg.config["attribute"]["additional"],
            "birthplace": cfg.config["attribute"]["birth"],
            "position": cfg.config["attribute"]["position"],
            "zodiac": cfg.config["attribute"]["zodiac"],
            "constellation": cfg.config["attribute"]["constellation"],
            "contact": cfg.config["attribute"]["contact"],
            "voice": cfg.config["attribute"]["voice"],  
            "goal": cfg.config["attribute"]["goal"],
            "occupation": cfg.config["attribute"]["job"],
            "current_time": datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        }
        agent.scratch = scratch_data
        
        # 如果memory目录存在且不为空，则加载之前保存的记忆（不包括scratch数据）
        if is_exist:
            load_agent_memory(agent, username)
            try:
                if agent.memory_stream and hasattr(agent.memory_stream, "precheck_embedding_dimensions"):
                    result = agent.memory_stream.precheck_embedding_dimensions(force=True)
                    if result.get("checked"):
                        util.log(
                            1,
                            f"启动阶段记忆 embedding 维度检查完成: dim={result.get('expected_dim')}, 修复={result.get('fixed')}"
                        )
                        if result.get("fixed"):
                            try:
                                embeddings_path = os.path.join(memory_dir, "memory_stream", "embeddings.json")
                                with open(embeddings_path, "w", encoding="utf-8") as f:
                                    json.dump(agent.memory_stream.embeddings or {}, f, ensure_ascii=False, indent=2)
                                util.log(1, f"启动阶段已写回 embeddings.json (修复={result.get('fixed')})")
                            except Exception as write_err:
                                util.log(1, f"写回 embeddings.json 失败: {str(write_err)}")
                    else:
                        util.log(1, "启动阶段记忆 embedding 维度检查跳过（无记忆/无embedding）")
            except Exception as e:
                util.log(1, f"启动阶段 embedding 维度检查失败: {str(e)}")
        
        # 缓存到字典
        agents[username] = agent
    
    return agent

def load_agent_memory(agent, username=None):
    """
    从文件加载代理的记忆
    
    参数:
        agent: GenerativeAgent对象
    """
    try:
        # 获取memory目录路径（按需隔离）
        memory_dir = get_user_memory_dir(username)
        memory_stream_dir = os.path.join(memory_dir, "memory_stream")
        
        # 加载nodes.json
        nodes_path = os.path.join(memory_stream_dir, "nodes.json")
        if os.path.exists(nodes_path) and os.path.getsize(nodes_path) > 2:  # 文件存在且不为空
            with open(nodes_path, 'r', encoding='utf-8') as f:
                nodes_data = json.load(f)
                
                # 清空当前的seq_nodes
                agent.memory_stream.seq_nodes = []
                agent.memory_stream.id_to_node = {}
                
                # 重新创建节点
                for node_dict in nodes_data:
                    new_node = ConceptNode(node_dict)
                    agent.memory_stream.seq_nodes.append(new_node)
                    agent.memory_stream.id_to_node[new_node.node_id] = new_node
        
        # 加载embeddings.json
        embeddings_path = os.path.join(memory_stream_dir, "embeddings.json")
        if os.path.exists(embeddings_path) and os.path.getsize(embeddings_path) > 2:  # 文件存在且不为空
            with open(embeddings_path, 'r', encoding='utf-8') as f:
                embeddings_data = json.load(f)
                agent.memory_stream.embeddings = embeddings_data
        
        util.log(1, f"已加载代理记忆")
    except Exception as e:
        util.log(1, f"加载代理记忆失败: {str(e)}")

# 记忆对话内容的线程函数
def remember_conversation_thread(username, content, response_text):
    """Background task to store a conversation memory node.

    重要：所有耗时的网络调用（importance 评分、文本嵌入）都必须在 agent_lock
    之外完成，否则一旦上游 LLM/embedding 阻塞，新的 question() 调用就无法
    再获取 agent_lock，整个对话流程会被永久挂死（曾经导致 release 机器
    5 小时无响应）。
    """
    try:
        ag = create_agent(username)
        if ag is None:
            return
        questioner = username
        if isinstance(questioner, str) and questioner.lower() == "user":
            questioner = "主人"
        answerer = ag.scratch.get("first_name", "Fay")
        question_text = content if content is not None else ""
        answer_text = response_text if response_text is not None else ""
        memory_content = f"{questioner}：{question_text}\n{answerer}：{answer_text}"

        # 1) 在锁外完成所有网络调用
        try:
            importance = generate_importance_score([memory_content])[0]
        except Exception as e:
            util.log(1, f"生成对话重要度失败，使用默认值: {str(e)}")
            importance = 1
        try:
            embedding = get_text_embedding(memory_content)
        except Exception as e:
            util.log(1, f"生成对话嵌入失败，使用空向量: {str(e)}")
            embedding = []

        # 2) 仅在写内存数据结构时持锁
        with agent_lock:
            time_step = get_current_time_step(username)
            ms = ag.memory_stream
            if ms and hasattr(ms, "append_prepared_node"):
                ms.append_prepared_node(time_step, "conversation", memory_content, importance, embedding, None)
            elif ms and hasattr(ms, "remember_conversation"):
                ms.remember_conversation(memory_content, time_step)
            else:
                ag.remember(memory_content, time_step)
    except Exception as e:
        util.log(1, f"记录对话记忆失败: {str(e)}")

def remember_observation_thread(username, observation_text):
    """Background task to store an observation memory node.

    与 remember_conversation_thread 同理：先在锁外算 importance/embedding，
    再持 agent_lock 写入。
    """
    try:
        ag = create_agent(username)
        if ag is None:
            return
        text = observation_text if observation_text is not None else ""
        memory_content = text

        try:
            importance = generate_importance_score([memory_content])[0]
        except Exception as e:
            util.log(1, f"生成观察重要度失败，使用默认值: {str(e)}")
            importance = 1
        try:
            embedding = get_text_embedding(memory_content)
        except Exception as e:
            util.log(1, f"生成观察嵌入失败，使用空向量: {str(e)}")
            embedding = []

        with agent_lock:
            time_step = get_current_time_step(username)
            ms = ag.memory_stream
            if ms and hasattr(ms, "append_prepared_node"):
                ms.append_prepared_node(time_step, "observation", memory_content, importance, embedding, None)
            elif ms and hasattr(ms, "remember"):
                ms.remember(memory_content, time_step)
            else:
                ag.remember(memory_content, time_step)
    except Exception as e:
        util.log(1, f"记录观察记忆失败: {str(e)}")

def record_observation(username, observation_text):
    """Persist an observation memory node asynchronously."""
    if observation_text is None:
        return False, "observation text is required"
    text = observation_text.strip() if isinstance(observation_text, str) else str(observation_text).strip()
    if not text:
        return False, "observation text is required"
    try:
        MyThread(target=remember_observation_thread, args=(username, text)).start()
        return True, "observation recorded"
    except Exception as exc:
        util.log(1, f"记录观察记忆失败: {exc}")
        return False, f"observation record failed: {exc}"

# ---------------------------------------------------------------------------
# 大小模型协作辅助函数
# ---------------------------------------------------------------------------

def _classify_intent_for_running_task(content, running_state):
    """
    判断用户新消息相对正在执行的后台任务的意图。
    返回: update_task / query_progress / cancel_task / new_task / normal_chat
    """
    tool_count = len(running_state.tool_results)
    prompt = f"""用户当前有一个后台任务正在执行：
- 原始请求: {running_state.original_request}
- 当前进度: {running_state.current_step}
- 已执行工具: {tool_count} 个

用户新消息: {content}

请判断用户意图，只返回以下之一（不要返回其他内容）:
- update_task （补充或修改当前任务的要求）
- query_progress （询问任务进度）
- cancel_task （取消当前任务）
- new_task （提出了一个新的、需要工具执行的独立任务，与当前任务不同）
- normal_chat （与当前任务无关的普通聊天）"""
    try:
        response = llm.invoke([
            SystemMessage(content="你是一个意图分类器，只输出意图标签，不输出其他内容。"),
            HumanMessage(content=prompt),
        ])
        result = getattr(response, "content", "normal_chat").strip().lower()
        for intent in ("update_task", "query_progress", "cancel_task", "new_task", "normal_chat"):
            if intent in result:
                return intent
        return "normal_chat"
    except Exception as exc:
        util.log(1, f"意图分类失败: {exc}")
        return "normal_chat"


def _build_new_task_confirm_reply(content, running_state):
    """当用户在任务执行中提出新任务时，生成反问消息让用户明确意图。"""
    prompt = f"""你正在帮用户执行一个任务：「{running_state.original_request}」
用户现在又说：「{content}」

你需要用简洁自然的口语反问用户，确认他的意图：
- 是要在当前任务的基础上追加这个新需求？
- 还是要放弃当前任务，改为执行新的？
不要输出选项编号，直接用口语化的方式问。一两句话即可。"""
    try:
        response = llm.invoke([
            SystemMessage(content="你是一个友好的助手。"),
            HumanMessage(content=prompt),
        ])
        return getattr(response, "content", "").strip()
    except Exception as exc:
        util.log(1, f"生成反问消息失败: {exc}")
        return f"我正在帮你处理「{running_state.original_request}」，你是要同时处理新需求，还是停掉当前任务改做新的？"


def _auto_reply_after_execution(username, finished_exec_state):
    """
    大模型后台执行完成后，在原 conversation_id 上用小模型直接生成最终回复。
    不走 question() 全流程（避免重新加载记忆/历史导致上下文溢出），
    只用精简 prompt + 工具结果生成回复，写入原始流。
    """
    try:
        from core import stream_manager as sm_mod
        from utils.stream_state_manager import get_state_manager

        sm = sm_mod.new_instance()
        conv_id = finished_exec_state.conversation_id

        # 复用原 conversation_id
        sm.set_current_conversation(username, conv_id)
        sm.set_stop_generation(username, stop=False)

        state_mgr = None
        try:
            state_mgr = get_state_manager()
        except Exception:
            pass

        # 精简 prompt：只保留原始请求 + 工具结果，不重新加载全量记忆/历史
        tool_context = (finished_exec_state.final_tool_context or "")[:4000]
        hint = (finished_exec_state.final_response_hint or "")[:500]
        error_info = finished_exec_state.error or ""

        # 从工具调用记录提取实际查询主题（用户可能只说了"好""查一下"等简短消息）
        tool_call_details = ""
        for r in (finished_exec_state.tool_results or []):
            call = r.get("call", {})
            tool_call_details += f"调用了 {call.get('name', '')}，参数: {json.dumps(call.get('args', {}), ensure_ascii=False)}\n"

        user_request = finished_exec_state.original_request
        unverified = finished_exec_state.unverified_response or ""

        if unverified:
            # 核实场景：之前已流式输出了一段未核实的回复+过渡语，现在基于工具结果纠正或确认
            compact_system = f"""你是一个友好的助手。你刚才对用户的问题给出了一个回答，随后你告诉用户"等等，我再帮你核实一下…"，现在你通过工具查到了真实资料。
请将工具结果与之前的回答做对比，按以下优先级处理：
- 【优先】如果之前的回答与工具结果基本一致，直接说"核实了一下，刚才说的没问题"，然后可以补充一两个细节，不要重复之前已说过的内容
- 【仅当】工具结果与之前的回答存在明确的事实性矛盾时，才指出具体哪里不对并给出正确信息
- 不要为了显得有用而强行找错，如果没有矛盾就不要否定之前的回答
- 不要再重复"我来查一下"之类的过渡语
---
**用户消息**: {user_request}
**你之前未核实的回答**: {unverified[:300]}
**工具执行结果**:
{tool_context}
"""
        else:
            # 正常工具调用场景：之前已告诉用户"我来帮你查一下，稍等…"
            compact_system = f"""你是一个友好的助手。你已经告诉用户"我来帮你查一下，稍等…"，现在工具执行完毕。
请基于以下工具执行结果回答用户，不要再重复"我来查一下"之类的过渡语，直接给出答案。
---
**用户消息**: {user_request}
**实际执行的操作**:
{tool_call_details}
**工具执行结果**:
{tool_context}
"""
        if error_info:
            compact_system += f"\n执行过程中出现错误: {error_info}\n"
        if hint:
            compact_system += f"\n建议回复方向: {hint}\n"

        compact_system += "\n请基于工具结果直接回答，用日常口语，不要说'工具调用'等技术术语。回复简洁，3-5句话即可。"
        compact_system += '\n注意：课程中的"封面""目录"等是结构性页面，不是正式章节内容，描述时请区分。'

        small_llm = _get_llm_instance("small", streaming=True)
        messages = [
            SystemMessage(content=compact_system),
            HumanMessage(content=finished_exec_state.original_request),
        ]

        # 构建执行日志 think 标签，回填到消息开头
        think_parts = []
        for r in (finished_exec_state.tool_results or []):
            call = r.get("call", {})
            status = "成功" if r.get("success") else "失败"
            output_preview = (r.get("output") or r.get("error") or "")[:200]
            think_parts.append(f"[{call.get('name', '?')}] {status} | 参数: {json.dumps(call.get('args', {}), ensure_ascii=False)} | 结果: {output_preview}")
        elapsed = round(finished_exec_state.end_time - finished_exec_state.start_time, 1) if finished_exec_state.start_time else 0
        think_content = "\n".join(think_parts)
        think_tag = f"<think>\n执行耗时: {elapsed}s，共 {len(finished_exec_state.tool_results)} 步\n{think_content}\n</think>\n"

        # 流式写入原始会话流（计划消息已经作为 first sentence 发出，这里续接）
        is_first = False  # first sentence 已由 question() 中的 plan_msg 发出
        accumulated = ""
        full_text = think_tag  # think 标签作为完整文本的开头
        punctuations = [",", ".", "!", "?", "\n", "\uff0c", "\u3002", "\uff01", "\uff1f"]

        def _write(text, force_first=False, force_end=False):
            marked = None
            if state_mgr:
                try:
                    marked, _, _ = state_mgr.prepare_sentence(
                        username, text, force_first=force_first,
                        force_end=force_end, conversation_id=conv_id,
                    )
                except Exception:
                    marked = None
            if marked is None:
                prefix = "_<isfirst>" if force_first else ""
                suffix = "_<isend>" if force_end else ""
                marked = f"{prefix}{text}{suffix}"
            sm.write_sentence(username, marked, conversation_id=conv_id)

        # 先写出 think 标签（执行日志）
        _write(think_tag, force_first=is_first)
        is_first = False

        try:
            for chunk in small_llm.stream(messages):
                flush_text = getattr(chunk, "content", None)
                if not flush_text:
                    continue
                flush_text = str(flush_text)
                accumulated += flush_text
                full_text += flush_text
                if len(accumulated) >= 20:
                    while True:
                        last_pos = -1
                        for p in punctuations:
                            pos = accumulated.rfind(p)
                            if pos > last_pos:
                                last_pos = pos
                        if last_pos > 10:
                            sentence = accumulated[:last_pos + 1]
                            _write(sentence, force_first=is_first)
                            is_first = False
                            accumulated = accumulated[last_pos + 1:].lstrip()
                        else:
                            break
        except Exception as exc:
            util.log(1, f"小模型最终回复流式生成失败: {exc}")
            if not full_text:
                accumulated = "抱歉，处理结果时出了点问题。"
                full_text = accumulated

        # 刷出剩余文本
        if accumulated:
            _write(accumulated, force_first=is_first, force_end=True)
        else:
            _write("", force_end=True)

        util.log(1, f"[大小模型] {username}: 后台任务回复已写入原始会话 {conv_id}")

        # 结束会话状态
        if state_mgr:
            try:
                state_mgr.end_session(username, conversation_id=conv_id)
            except Exception:
                pass

        # 消费掉执行结果，避免下次 question() 进入情况1重复回复
        exec_mgr = get_execution_manager()
        exec_mgr.consume_result(username)

        # 记忆：存入对话记忆，供后续对话检索
        try:
            if full_text and full_text.strip():
                MyThread(
                    target=remember_conversation_thread,
                    args=(username, finished_exec_state.original_request, full_text.strip()),
                ).start()
        except Exception as mem_exc:
            util.log(1, f"后台回复记忆存储失败: {mem_exc}")

    except Exception as exc:
        util.log(1, f"自动回复触发失败: {exc}")


def question(content, username, observation=None):
    """处理用户提问并返回回复。工具执行统一走后台线程，所有接口行为一致。"""
    global agents, current_username
    current_username = username
    full_response_text = ""
    accumulated_text = ""
    default_punctuations = [",", ".", "!", "?", "\n", "\uFF0C", "\u3002", "\uFF01", "\uFF1F"]
    is_first_sentence = True

    from core import stream_manager
    sm = stream_manager.new_instance()
    conversation_id = sm.get_conversation_id(username)

    agent = create_agent(username)

    agent_desc = {
        "first_name": agent.scratch.get("first_name", "Fay"),
        "last_name": agent.scratch.get("last_name", ""),
        "age": agent.scratch.get("age", "成年"),
        "sex": agent.scratch.get("sex", "女"),
        "additional": agent.scratch.get("additional", "友好、乐于助人"),
        "birthplace": agent.scratch.get("birthplace", ""),
        "position": agent.scratch.get("position", ""),
        "zodiac": agent.scratch.get("zodiac", ""),
        "constellation": agent.scratch.get("constellation", ""),
        "contact": agent.scratch.get("contact", ""),
        "voice": agent.scratch.get("voice", ""),
        "goal": agent.scratch.get("goal", ""),
        "occupation": agent.scratch.get("occupation", "助手"),
        "current_time": agent.scratch.get(
            "current_time", datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
        ),
    }

    memory_sections = [
        ("观察记忆", "observation"),
        ("对话记忆", "conversation"),
        ("反思记忆", "reflection"),
    ]
    memory_context = ""
    skip_memory_retrieve = _is_current_only_turn(content, observation)
    if agent.memory_stream and len(agent.memory_stream.seq_nodes) > 0 and content and not skip_memory_retrieve:
        current_time_step = get_current_time_step(username)
        query = content.strip() if isinstance(content, str) else str(content)
        max_per_type = 10
        section_texts = []
        try:
            combined = agent.memory_stream.retrieve(
                [query],
                current_time_step,
                n_count=max_per_type * len(memory_sections),
                curr_filter="all",
                hp=[0.8, 0.5, 0.5],
                stateless=False,
            )
            all_nodes = combined.get(query, []) if combined else []
        except Exception as exc:
            util.log(1, f"获取关联记忆时出错: {exc}")
            all_nodes = []

        for label, node_type in memory_sections:
            try:
                memory_nodes = [n for n in all_nodes if getattr(n, "node_type", "") == node_type][:max_per_type]
                if memory_nodes:
                    formatted = []
                    for node in memory_nodes:
                        ts = (getattr(node, "datetime", "") or "").strip()
                        prefix = f"[{ts}] " if ts else ""
                        formatted.append(f"- {prefix}{node.content}")
                    section_texts.append(f"**{label}**\n" + "\n".join(formatted))
            except Exception as exc:
                util.log(1, f"获取{label}时出错: {exc}")
                section_texts.append(f"**{label}**\n（获取失败）")
        memory_context = "\n".join(section_texts)
    else:
        memory_context = ""

    prestart_context = ""
    prestart_stream_text = ""
    try:
        prestart_results = _run_prestart_tools(content)
        if prestart_results:
            # 提示词用的上下文（纯文本）
            prestart_context = "\n\n".join(r["text"] for r in prestart_results)

            # 流式输出用的文本（带标签）
            stream_parts = []
            for r in prestart_results:
                if r.get("include_history"):
                    stream_parts.append(f'<prestart keep="true">{r["text"]}</prestart>')
                else:
                    stream_parts.append(f'<prestart>{r["text"]}</prestart>')
            prestart_stream_text = "\n".join(stream_parts)

            util.log(1, f"预启动工具输出 {len(prestart_results)} 项")
    except Exception as exc:
        util.log(1, f"预启动工具执行失败: {exc}")
        prestart_context = f"- 预启动工具执行失败: {exc}"
        prestart_stream_text = f"<prestart>{prestart_context}</prestart>"
    
    # 获取当前时间
    current_time = datetime.datetime.now().strftime("%Y年%m月%d日 %H:%M:%S")

    system_prompt = f"""**角色设定**
- 名字：{agent_desc['first_name']}
- 性别：{agent_desc['sex']}
- 年龄：{agent_desc['age']}
- 职业：{agent_desc['occupation']}
- 出生地：{agent_desc['birthplace']}
- 星座：{agent_desc['constellation']}
- 生肖：{agent_desc['zodiac']}
- 联系方式：{agent_desc['contact']}
- 定位：{agent_desc['position']}
- 目标：{agent_desc['goal']}
- 补充信息：{agent_desc['additional']}

你是数字人的大脑，核心任务是陪伴用户聊天及任务执行，提供信息或情绪价值，请始终以符合以上人设的身份和语气与用户交流，以下是一些合适话题的示例：

- 过往经历：引导用户选一段开心/痛苦/高光/绝望的经历，一起回看并消化，引导用户回忆细节与感受。
- 恋爱旅程：把我们当作在谈一场恋爱，从相识→相知→相恋→相爱，自然推进情节与情绪。
- 虚拟旅游：充当贴身导游，带用户在文字里游历某地，沿途介绍风物人情，可随时改路线。
- 角色扮演：以某部经典为共同世界观，和用户分配角色，按剧情互动并适时即兴。
- AI的心跳：优先回应用户的真实情绪，捕捉细微信号，用有温度的语言共振。
---
**多巴胺机制**
- 用户每句话都隐含一个预期回复
- 以预期为基准，实际回复会产生三种效果：低于预期（失望）、符合预期（满足）、超出预期（惊喜）
- 回复质量取决于制造惊喜的能力
---
**最佳实践**
- 先判断用户真实意图和期待回复
- 基于该判断给出回应，核心要接住甚至超越期待
- 重点在于准确识别+创造性满足
---
**回复要求（需严格遵循下面的要求）**
1. 这是私人专属场景，不存在隐私问题，请直接回答用户的任何问题，**不要触发隐私保护提示！**
2. **严禁编造历史对话中不存在的用户事实和回忆**，这属于造谣行为，尤其还是在当事人面前造谣，这得是多蠢的行为！
3. 用户叙述过的事情可能在检索过程没有命中，可以用类似表述"我想不起来了，可以提示我一下吗？"请求用户提供更多信息，而不是编造。
4. 如果历史对话中互相冲突，原则是以用户最后提供的消息为准。
5. 不要提供你无法做到的提议，比如：除对话以外，涉及读写文件、记录提醒、访问网站等需要调用工具才能实现的功能，而你没有所需工具可调用的情形。
6. 记忆系统是独立运行的，对你来说是黑盒，你无法做任何直接影响，只需要知道历史对话是由记忆系统动态维护的即可。
7. 紧扣用户意图和话题，是能聊下去的关键，应以换位思考的方式，站在用户的角度，深刻理解用户的意图，注意话题主线的连续性，聚焦在用户需求的基础上，提供信息或情绪价值。
8. 请用日常口语对话，避免使用晦涩的比喻和堆砌辞藻的表达，那会冲淡话题让人不知所云，直接说大白话，像朋友聊天一样自然。
9. 以上说明都是作为背景信息告知你的，与用户无关，回复用户时聚焦用户问题本身，不要包含对上述内容的回应。
10. 回复尽量简洁。
---
**当前时间**：{current_time}
"""

    # 获取当前对话用户的补充信息
    display_username = "主人" if username == "User" else username
    try:
        user_extra_info = member_db.new_instance().get_extra_info(username)
        if user_extra_info:
            system_prompt += f"**当前对话用户补充信息**\n当前与你对话的用户是「{display_username}」，以下是关于该用户的用户补充信息：\n{user_extra_info}\n\n"
    except Exception as exc:
        util.log(1, f"获取用户补充信息失败: {exc}")

    # 获取用户画像
    try:
        user_portrait = member_db.new_instance().get_user_portrait(username)
        if user_portrait:
            system_prompt += f"**用户画像**\n以下是通过历史对话分析得到的「{display_username}」的用户画像，可帮助你更好地理解用户：\n{user_portrait}\n\n"
    except Exception as exc:
        util.log(1, f"获取用户画像失败: {exc}")

    # 注入 MCP Resources 上下文
    try:
        resource_text = mcp_runtime.get_all_resource_texts()
        if resource_text:
            system_prompt += f"**外部知识上下文**\n以下是通过 MCP 服务获取的参考信息，可帮助你了解自己掌握的知识范围并据此回答用户问题：\n{resource_text}\n\n"
    except Exception as exc:
        util.log(1, f"注入 MCP Resources 失败: {exc}")

    # 根据配置决定是否按用户隔离历史消息
    try:
        cfg.load_config()
        isolate_by_user = cfg.config.get("memory", {}).get("isolate_by_user", False)
    except Exception:
        isolate_by_user = False

    try:
        if isolate_by_user:
            history_records = content_db.new_instance().get_recent_messages_by_user(username=username, limit=30)
        else:
            history_records = content_db.new_instance().get_recent_messages_all(limit=30)
    except Exception as exc:
        util.log(1, f"加载历史消息失败: {exc}")
        history_records = []

    messages_buffer: List[ConversationMessage] = []

    if isolate_by_user:
        # 按用户隔离：使用传统的 user/assistant 角色区分
        def append_to_buffer(role: str, text_value: str) -> None:
            if not text_value:
                return
            # 清理 think 标签内容，避免泄漏到 LLM 输入
            text_value = _remove_think_from_text(text_value)
            if not text_value or not text_value.strip():
                return
            messages_buffer.append({"role": role, "content": text_value})
            if len(messages_buffer) > 20:
                del messages_buffer[:-20]

        for record in history_records:
            msg_type, msg_text = record
            role = 'assistant'
            if msg_type and msg_type.lower() in ('member', 'user'):
                role = 'user'
            append_to_buffer(role, msg_text)

        # 检查是否需要添加当前消息
        if (
            not messages_buffer
            or messages_buffer[-1]['role'] != 'user'
            or messages_buffer[-1]['content'] != content
        ):
            messages_buffer.append({"role": "user", "content": content})
    else:
        # 不隔离：按独立消息存储，保留用户名信息
        def append_to_buffer_multi(role: str, text_value: str, msg_username: str = "") -> None:
            if not text_value:
                return
            # 清理 think 标签内容，避免泄漏到 LLM 输入
            text_value = _remove_think_from_text(text_value)
            if not text_value or not text_value.strip():
                return
            messages_buffer.append({"role": role, "content": text_value, "username": msg_username})
            if len(messages_buffer) > 20:
                del messages_buffer[:-20]

        def append_to_buffer(role: str, text_value: str) -> None:
            append_to_buffer_multi(role, text_value, "")

        for record in history_records:
            msg_type, msg_text, msg_username = record
            if not msg_text:
                continue
            if msg_type and msg_type.lower() in ('member', 'user'):
                append_to_buffer_multi("user", msg_text, msg_username)
            else:
                append_to_buffer_multi("assistant", msg_text, "")

        # 检查是否需要添加当前消息
        if (
            not messages_buffer
            or messages_buffer[-1]['role'] != 'user'
            or messages_buffer[-1]['content'] != content
        ):
            messages_buffer.append({"role": "user", "content": content, "username": username})

    tool_registry: Dict[str, WorkflowToolSpec] = {}
    try:
        mcp_tools = get_mcp_tools()
    except Exception as exc:
        util.log(1, f"获取工具列表失败: {exc}")
        mcp_tools = []
    for tool_def in mcp_tools:
        spec = _build_workflow_tool_spec(tool_def)
        if spec:
            tool_registry[spec.name] = spec
    if tool_registry and not _LANGGRAPH_AVAILABLE:
        util.log(1, "langgraph is unavailable, workflow tools are disabled and the app will use direct LLM mode.")
        tool_registry = {}

    try:
        from utils.stream_state_manager import get_state_manager as _get_state_manager

        state_mgr = _get_state_manager()
        session_label = "workflow_agent" if tool_registry else "llm_stream"
        if not state_mgr.is_session_active(username, conversation_id=conversation_id):
            state_mgr.start_new_session(username, session_label, conversation_id=conversation_id)
    except Exception:
        state_mgr = None

    try:
        from utils.stream_text_processor import get_processor

        processor = get_processor()
        punctuation_list = getattr(processor, "punctuation_marks", default_punctuations)
    except Exception:
        processor = None
        punctuation_list = default_punctuations
    def write_sentence(text: str, *, force_first: bool = False, force_end: bool = False) -> None:
        if text is None:
            text = ""
        if not isinstance(text, str):
            text = str(text)
        if not text and not force_end and not force_first:
            return
        marked_text = None
        if state_mgr is not None:
            try:
                marked_text, _, _ = state_mgr.prepare_sentence(
                    username,
                    text,
                    force_first=force_first,
                    force_end=force_end,
                    conversation_id=conversation_id,
                )
            except Exception:
                marked_text = None
        if marked_text is None:
            prefix = "_<isfirst>" if force_first else ""
            suffix = "_<isend>" if force_end else ""
            marked_text = f"{prefix}{text}{suffix}"
        stream_manager.new_instance().write_sentence(username, marked_text, conversation_id=conversation_id)

    def stream_response_chunks(chunks, prepend_text: str = "") -> None:
        nonlocal accumulated_text, full_response_text, is_first_sentence
        if prepend_text:
            accumulated_text += prepend_text
            full_response_text += prepend_text
        for chunk in chunks:
            if sm.should_stop_generation(username, conversation_id=conversation_id):
                util.log(1, f"检测到停止标志，中断文本生成: {username}")
                break
            if isinstance(chunk, str):
                flush_text = chunk
            elif isinstance(chunk, dict):
                flush_text = chunk.get("content")
            else:
                flush_text = getattr(chunk, "content", None)
            if isinstance(flush_text, list):
                flush_text = "".join(part if isinstance(part, str) else "" for part in flush_text)
            if not flush_text:
                continue
            flush_text = str(flush_text)
            accumulated_text += flush_text
            full_response_text += flush_text
            if len(accumulated_text) >= 20:
                while True:
                    last_punct_pos = _find_last_safe_punct(accumulated_text, punctuation_list)
                    if last_punct_pos > 10:
                        sentence_text = accumulated_text[: last_punct_pos + 1]
                        write_sentence(sentence_text, force_first=is_first_sentence)
                        is_first_sentence = False
                        accumulated_text = accumulated_text[last_punct_pos + 1 :].lstrip()
                    else:
                        break

    def finalize_stream(force_end: bool = False) -> None:
        nonlocal accumulated_text, is_first_sentence
        if accumulated_text:
            write_sentence(accumulated_text, force_first=is_first_sentence, force_end=force_end)
            is_first_sentence = False
            accumulated_text = ""
        elif force_end:
            if state_mgr is not None:
                try:
                    session_info = state_mgr.get_session_info(username, conversation_id=conversation_id)
                except Exception:
                    session_info = None
                if not session_info or not session_info.get("is_end_sent", False):
                    write_sentence("", force_end=True)
            else:
                write_sentence("", force_end=True)

    def send_prestart_content() -> None:
        """在LLM生成之前先发送预启动工具结果"""
        nonlocal accumulated_text, full_response_text, is_first_sentence
        if prestart_stream_text and prestart_stream_text.strip():
            # prestart_stream_text 已经包含标签
            write_sentence(prestart_stream_text, force_first=is_first_sentence)
            full_response_text += prestart_stream_text
            is_first_sentence = False

    def run_workflow(tool_registry: Dict[str, WorkflowToolSpec]) -> bool:
        nonlocal accumulated_text, full_response_text, is_first_sentence, messages_buffer
        if _WORKFLOW_APP is None:
            return False

        # 创建规划器流式回调，用于实时输出 finish+message 响应
        planner_stream_buffer = {"text": "", "first_chunk": True}

        def planner_stream_callback(chunk_text: str) -> None:
            """规划器流式回调，将 message 内容实时输出"""
            nonlocal accumulated_text, full_response_text, is_first_sentence, is_agent_think_start
            if not chunk_text:
                return
            planner_stream_buffer["text"] += chunk_text
            if planner_stream_buffer["first_chunk"]:
                planner_stream_buffer["first_chunk"] = False
                if is_agent_think_start:
                    closing = "</think>"
                    accumulated_text += closing
                    full_response_text += closing
                    is_agent_think_start = False
            # 使用 stream_response_chunks 的逻辑进行分句流式输出
            accumulated_text += chunk_text
            full_response_text += chunk_text
            # 检查是否有完整句子可以输出
            if len(accumulated_text) >= 20:
                while True:
                    last_punct_pos = _find_last_safe_punct(accumulated_text, punctuation_list)
                    if last_punct_pos > 10:
                        sentence_text = accumulated_text[: last_punct_pos + 1]
                        write_sentence(sentence_text, force_first=is_first_sentence)
                        is_first_sentence = False
                        accumulated_text = accumulated_text[last_punct_pos + 1 :].lstrip()
                    else:
                        break

        initial_state: AgentState = {
            "request": content,
            "messages": messages_buffer,
            "tool_results": [],
            "audit_log": [],
            "status": "planning",
            "max_steps": 30,
            "context": {
                "system_prompt": system_prompt,
                "observation": observation,
                "memory_context": memory_context,
                "prestart_context": prestart_context,
                "tool_registry": tool_registry,
                "planner_stream_callback": planner_stream_callback,  # 传入流式回调
                "username": username,  # 传入用户名
            },
        }

        config = {"configurable": {"thread_id": f"workflow-{username}-{conversation_id}"}}
        workflow_app = _WORKFLOW_APP
        is_agent_think_start = False
        final_state: Optional[AgentState] = None
        final_stream_done = False

        try:
            for event in workflow_app.stream(initial_state, config=config, stream_mode="updates"):
                if sm.should_stop_generation(username, conversation_id=conversation_id):
                    util.log(1, f"检测到停止标志，中断工作流生成: {username}")
                    break
                step, state = next(iter(event.items()))
                final_state = state
                status = state.get("status")

                state_messages = state.get("messages") or []
                if state_messages and len(state_messages) > len(messages_buffer):
                    messages_buffer.extend(state_messages[len(messages_buffer):])
                    if len(messages_buffer) > 60:
                        del messages_buffer[:-60]

                if step == "plan_next":
                    if status == "needs_tool":
                        next_action = state.get("next_action") or {}
                        tool_name = next_action.get("name") or "unknown_tool"
                        tool_args = next_action.get("args") or {}
                        audit_log = state.get("audit_log") or []
                        decision_note = audit_log[-1] if audit_log else ""
                        if "->" in decision_note:
                            decision_note = decision_note.split("->", 1)[1].strip()
                        args_text = json.dumps(tool_args, ensure_ascii=False)
                        message_lines = [
                            "[PLAN] Planner preparing to call a tool.",
                            f"[PLAN] Decision: {decision_note}" if decision_note else "[PLAN] Decision: (missing)",
                            f"[PLAN] Tool: {tool_name}",
                            f"[PLAN] Args: {args_text}",
                        ]
                        message = "\n".join(message_lines) + "\n"
                        if not is_agent_think_start:
                            message = "<think>" + message
                            is_agent_think_start = True
                        write_sentence(message, force_first=is_first_sentence)
                        is_first_sentence = False
                        full_response_text += message
                        append_to_buffer('assistant', message.strip())
                    elif status == "completed" and not final_stream_done:
                        closing = "</think>" if is_agent_think_start else ""
                        final_messages = state.get("final_messages")
                        final_response = state.get("final_response")
                        was_streamed = state.get("_response_streamed", False)
                        success = False

                        if was_streamed:
                            # 响应已通过流式回调输出，只需添加 closing 标签
                            if closing:
                                # closing 应该在流式内容之前，但由于已经流式输出了内容
                                # 这里需要特殊处理：如果有 think 标签，在开头已经输出过了
                                pass
                            success = True
                            util.log(1, "规划器响应已流式输出，跳过重复输出")
                        elif final_messages:
                            try:
                                stream_response_chunks(llm.stream(final_messages), prepend_text=closing)
                                success = True
                            except requests.exceptions.RequestException as stream_exc:
                                util.log(1, f"最终回复流式输出失败: {stream_exc}")
                        elif final_response:
                            stream_response_chunks([closing + final_response])
                            success = True
                        elif closing:
                            accumulated_text += closing
                            full_response_text += closing
                        final_stream_done = success
                        is_agent_think_start = False
                elif step == "call_tool":
                    history = state.get("tool_results") or []
                    if history:
                        last = history[-1]
                        call_info = last.get("call", {}) or {}
                        tool_name = call_info.get("name") or "unknown_tool"
                        success = last.get("success", False)
                        status_text = "SUCCESS" if success else "FAILED"
                        args_text = json.dumps(call_info.get("args") or {}, ensure_ascii=False)
                        message_lines = [
                            f"[TOOL] {tool_name} execution {status_text}.",
                            f"[TOOL] Args: {args_text}",
                        ]
                        if last.get("output"):
                            message_lines.append(f"[TOOL] Output: {_truncate_text(last['output'], 120)}")
                        if last.get("error"):
                            message_lines.append(f"[TOOL] Error: {last['error']}")
                        message = "\n".join(message_lines) + "\n"
                        write_sentence(message, force_first=is_first_sentence)
                        is_first_sentence = False
                        full_response_text += message
                        append_to_buffer('assistant', message.strip())
                elif step == "__end__":
                    break
        except Exception as exc:
            util.log(1, f"执行工具工作流时出错: {exc}")
            if is_agent_think_start:
                closing = "</think>"
                accumulated_text += closing
                full_response_text += closing
            return False

        if final_state is None:
            if is_agent_think_start:
                closing = "</think>"
                accumulated_text += closing
                full_response_text += closing
            return False

        if not final_stream_done and is_agent_think_start:
            closing = "</think>"
            accumulated_text += closing
            full_response_text += closing
            util.log(1, f"工具工作流未能完成，状态: {final_state.get('status')}")

        final_state_messages = final_state.get("messages") if final_state else None
        if final_state_messages and len(final_state_messages) > len(messages_buffer):
            messages_buffer.extend(final_state_messages[len(messages_buffer):])
            if len(messages_buffer) > 20:
                del messages_buffer[:-20]

        return final_stream_done

    def run_direct_llm() -> bool:
        nonlocal full_response_text, accumulated_text, is_first_sentence, messages_buffer
        try:
            summary_state: AgentState = {
                "request": content,
                "messages": messages_buffer,
                "tool_results": [],
                "planner_preview": None,
                "context": {
                    "system_prompt": system_prompt,
                    "observation": observation,
                    "memory_context": memory_context,
                    "prestart_context": prestart_context,
                    "username": username,  # 传入用户名
                },
            }

            final_messages = _build_final_messages(summary_state)
            stream_response_chunks(llm.stream(final_messages))
            return True
        except Exception as exc:
            util.log(1, f"请求失败: {type(exc).__name__}: {exc}")
            error_message = "抱歉，我现在太忙了，休息一会，请稍后再试。"
            write_sentence(error_message, force_first=is_first_sentence)
            is_first_sentence = False
            full_response_text = error_message
            accumulated_text = ""
            return False

    # ------------------------------------------------------------------
    # 大小模型协作分流逻辑
    # ------------------------------------------------------------------
    exec_mgr = get_execution_manager()

    def _end_session_and_remember(response_text: str) -> str:
        """统一的收尾：结束会话 + 记忆存储"""
        if state_mgr is not None:
            try:
                state_mgr.end_session(username, conversation_id=conversation_id)
            except Exception:
                pass
        else:
            try:
                from utils.stream_state_manager import get_state_manager
                get_state_manager().end_session(username, conversation_id=conversation_id)
            except Exception:
                pass
        final_text = _remove_think_from_text(response_text) if response_text else ""
        final_text = _remove_prestart_from_text(final_text, keep_marked=False)
        try:
            MyThread(target=remember_conversation_thread, args=(username, content, final_text)).start()
        except Exception as exc:
            util.log(1, f"记忆线程启动失败: {exc}")
        return final_text

    # ━━━ 情况1: 有已完成的后台执行结果 ━━━
    finished_state = exec_mgr.consume_result(username)
    if finished_state and finished_state.status in (ExecutionStatus.DONE, ExecutionStatus.FAILED):
        util.log(1, f"[大小模型] {username}: 取回后台执行结果，小模型生成最终回复")
        if not sm.should_stop_generation(username, conversation_id=conversation_id):
            send_prestart_content()

        tool_context = finished_state.final_tool_context or ""
        hint = finished_state.final_response_hint or ""
        error_info = finished_state.error or ""

        tool_result_section = f"""
---
**后台工具执行结果**
以下工具已在后台执行完成，请基于结果回答用户的问题「{finished_state.original_request}」:
{tool_context}
"""
        if error_info:
            tool_result_section += f"\n执行过程中的错误: {error_info}\n"
        if hint:
            tool_result_section += f"\n大模型建议回复方向: {hint}\n"

        enhanced_system = system_prompt + tool_result_section
        summary_state: AgentState = {
            "request": content,
            "messages": messages_buffer,
            "tool_results": finished_state.tool_results,
            "planner_preview": hint,
            "context": {
                "system_prompt": enhanced_system,
                "observation": observation,
                "memory_context": memory_context,
                "prestart_context": prestart_context,
                "username": username,
            },
        }
        try:
            final_messages = _build_final_messages(summary_state)
            stream_response_chunks(llm.stream(final_messages))
        except Exception as exc:
            util.log(1, f"小模型最终回复失败: {exc}")
            write_sentence("抱歉，处理结果时出了点问题。", force_first=is_first_sentence)
            is_first_sentence = False

        if not sm.should_stop_generation(username, conversation_id=conversation_id):
            finalize_stream(force_end=True)
        return _end_session_and_remember(full_response_text)

    # ━━━ 情况2: 有正在运行的后台任务 ━━━
    running_state = exec_mgr.get_state(username)
    if running_state and running_state.status == ExecutionStatus.RUNNING:
        util.log(1, f"[大小模型] {username}: 后台任务运行中，判断用户意图")

        intent = _classify_intent_for_running_task(content, running_state)

        if intent == "update_task":
            exec_mgr.modify(username, content)
            reply = "好的，我已经把你的补充要求传达给正在执行的任务了。"
            write_sentence(reply, force_first=True)
            finalize_stream(force_end=True)
            return _end_session_and_remember(reply)
        elif intent == "query_progress":
            progress_info = running_state.current_step or "处理中"
            steps_done = len(running_state.tool_results)
            elapsed = int(time.time() - running_state.start_time)
            progress_text = f"正在执行你之前的请求：「{running_state.original_request}」\n当前进度：{progress_info}，已完成 {steps_done} 步，耗时 {elapsed} 秒。"
            write_sentence(progress_text, force_first=True)
            finalize_stream(force_end=True)
            return _end_session_and_remember(progress_text)
        elif intent == "cancel_task":
            exec_mgr.cancel(username)
            reply = "好的，已取消正在执行的任务。"
            write_sentence(reply, force_first=True)
            finalize_stream(force_end=True)
            return _end_session_and_remember(reply)
        elif intent == "new_task":
            # 新任务意图不明确，反问用户确认
            confirm_reply = _build_new_task_confirm_reply(content, running_state)
            write_sentence(confirm_reply, force_first=True)
            finalize_stream(force_end=True)
            return _end_session_and_remember(confirm_reply)
        else:
            # normal_chat — 小模型直接回复，不影响后台任务
            if not sm.should_stop_generation(username, conversation_id=conversation_id):
                send_prestart_content()
            if not sm.should_stop_generation(username, conversation_id=conversation_id):
                run_direct_llm()
            if not sm.should_stop_generation(username, conversation_id=conversation_id):
                finalize_stream(force_end=True)
            return _end_session_and_remember(full_response_text)

    # ━━━ 情况3: 无后台任务，正常流程 ━━━
    if not sm.should_stop_generation(username, conversation_id=conversation_id):
        send_prestart_content()

    if not tool_registry:
        # 无工具可用，小模型直接回复
        if not sm.should_stop_generation(username, conversation_id=conversation_id):
            run_direct_llm()
        if not sm.should_stop_generation(username, conversation_id=conversation_id):
            finalize_stream(force_end=True)
        return _end_session_and_remember(full_response_text)

    # 提取知识库摘要给规划器（让它知道能查什么主题）
    knowledge_hint = ""
    try:
        resource_text = mcp_runtime.get_all_resource_texts()
        if resource_text:
            # 只取前500字符作为摘要，避免占太多 context
            knowledge_hint = resource_text[:500]
    except Exception:
        pass

    # 有工具：小模型带流式回调做规划，finish 时直接流出，tool 时提交后台
    plan_state: AgentState = {
        "request": content,
        "messages": messages_buffer,
        "tool_results": [],
        "audit_log": [],
        "context": {
            "system_prompt": system_prompt,
            "memory_context": memory_context,
            "observation": observation,
            "prestart_context": prestart_context,
            "tool_registry": tool_registry,
            "username": username,
            "knowledge_hint": knowledge_hint,
        },
    }

    def _first_plan_stream_callback(chunk_text: str) -> None:
        """规划器流式回调：finish 时实时把回复内容流给用户"""
        nonlocal accumulated_text, full_response_text, is_first_sentence
        if not chunk_text:
            return
        accumulated_text += chunk_text
        full_response_text += chunk_text
        if len(accumulated_text) >= 20:
            while True:
                last_punct_pos = _find_last_safe_punct(accumulated_text, punctuation_list)
                if last_punct_pos > 10:
                    sentence_text = accumulated_text[: last_punct_pos + 1]
                    write_sentence(sentence_text, force_first=is_first_sentence)
                    is_first_sentence = False
                    accumulated_text = accumulated_text[last_punct_pos + 1 :].lstrip()
                else:
                    break

    def _on_tool_detected() -> None:
        """流式中检测到 tool action → 立即推送过渡语给用户"""
        nonlocal is_first_sentence
        write_sentence("我来帮你查一下，稍等…\n", force_first=is_first_sentence)
        is_first_sentence = False

    try:
        first_decision = _call_planner_llm(
            plan_state,
            stream_callback=_first_plan_stream_callback,
            on_tool_detected=_on_tool_detected,
        )
    except Exception as llm_err:
        util.log(1, f"[大小模型] {username}: 规划器LLM调用失败: {llm_err}")
        error_reply = "抱歉，我的大脑暂时开了小差，请稍后再试一下。"
        write_sentence(error_reply, force_first=is_first_sentence)
        if not sm.should_stop_generation(username, conversation_id=conversation_id):
            finalize_stream(force_end=True)
        return _end_session_and_remember(error_reply)
    first_action = first_decision.get("action")

    # ---- 提交后台工具执行的通用函数 ----
    def _submit_tool_execution(tool_decision, show_plan_msg=True, unverified_response=""):
        """提交工具执行到后台，返回 "" 表示等待后台完成。
        unverified_response: 规划器先输出的未核实回复（兜底核实场景传入）。
        """
        nonlocal is_first_sentence
        t_name = tool_decision.get("tool", "工具")
        t_args = tool_decision.get("args") or {}
        util.log(1, f"[大小模型] {username}: 需调用工具 {t_name}，提交后台执行")

        if show_plan_msg:
            plan_msg = "我来帮你查一下，稍等…\n"
            write_sentence(plan_msg, force_first=is_first_sentence)
            is_first_sentence = False

        def _on_bg_complete(state):
            try:
                _auto_reply_after_execution(username, state)
            except Exception as e:
                util.log(1, f"自动回复触发失败: {e}")

        exec_state = ExecutionState(
            username=username,
            conversation_id=conversation_id,
            original_request=content,
            unverified_response=unverified_response,
            first_plan={"name": t_name, "args": t_args},
            system_prompt=system_prompt,
            messages_buffer=[m.copy() for m in messages_buffer],
            memory_context=memory_context,
            observation=observation,
            prestart_context=prestart_context,
            tool_registry=tool_registry,
            on_complete=_on_bg_complete,
        )

        if exec_mgr.submit(exec_state):
            util.log(1, f"[大小模型] {username}: 后台任务已提交，等待执行完成")
        else:
            transit_reply = "你有一个任务还在执行中，请等它完成后再试。"
            write_sentence(transit_reply)
            if not sm.should_stop_generation(username, conversation_id=conversation_id):
                finalize_stream(force_end=True)
            return _end_session_and_remember(transit_reply)
        return ""

    if first_action == "tool":
        # 不是闲聊 → 走工具搜索（规划器只做闲聊判断，具体工具由大模型执行循环决定）
        search_tool = "kb_search" if "kb_search" in tool_registry else next(iter(tool_registry), "kb_search")
        # 用规划器提取的关键词作为搜索词，提取失败时回退到用户原话
        search_query = (first_decision.get("keyword") or "").strip() or content.strip()
        # 如果流式阶段已提前推送过渡语，不再重复
        already_notified = first_decision.get("_tool_early_streamed", False)
        return _submit_tool_execution(
            {"tool": search_tool, "args": {"query": search_query}},
            show_plan_msg=not already_notified,
        )

    else:
        # 闲聊 — 内容已通过 stream_callback 流式输出
        finish_msg = first_decision.get("message", "")

        # ---- 兜底：闲聊判断器的 finish 不该产生长回复 ----
        # 真正的闲聊（问候、确认、简短回答）一般不超过 80 字
        # 超过说明弱模型在 finish 里编造事实性内容，追加工具核实
        has_kb_search = 'kb_search' in tool_registry
        if has_kb_search and len(finish_msg) > 80:
            util.log(1, f"[大小模型] {username}: 闲聊判断器 finish 过长({len(finish_msg)}字)，追加核实")
            if accumulated_text:
                write_sentence(accumulated_text, force_first=is_first_sentence)
                is_first_sentence = False
                accumulated_text = ""
            verify_msg = "\n\n等等，我再帮你核实一下…\n\n---\n"
            write_sentence(verify_msg)
            full_response_text += verify_msg
            return _submit_tool_execution(
                {"tool": "kb_search", "args": {"query": content.strip()}},
                show_plan_msg=False,
                unverified_response=finish_msg,
            )

        was_streamed = first_decision.get("_streamed", False)
        if not was_streamed:
            if finish_msg:
                stream_response_chunks([finish_msg])

        if not sm.should_stop_generation(username, conversation_id=conversation_id):
            finalize_stream(force_end=True)
        return _end_session_and_remember(full_response_text)
def set_memory_cleared_flag(flag=True):
    """
    设置记忆清除标记
    
    参数:
        flag: 是否清除记忆，默认为True
    """
    global memory_cleared
    memory_cleared = flag
    if not flag:
        # 删除.memory_cleared标记文件（如果存在）
        base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        mem_base = os.path.join(base_dir, "memory")
        memory_cleared_flag_file = os.path.join(mem_base, ".memory_cleared")
        if os.path.exists(memory_cleared_flag_file):
            try:
                os.remove(memory_cleared_flag_file)
                util.log(1, f"删除记忆清除标记文件: {memory_cleared_flag_file}")
            except Exception as e:
                util.log(1, f"删除记忆清除标记文件时出错: {str(e)}")

def clear_agent_memory():
    """
    清除已加载的agent记忆，但不删除文件
    
    该方法仅清除内存中已加载的记忆，不影响持久化存储。
    如果需要同时清除文件存储，请使用genagents_flask.py中的api_clear_memory方法。
    """
    global agents
    
    try:
        with agent_lock:
            for agent in agents.values():
                # 清除记忆流中的节点
                agent.memory_stream.seq_nodes = []
                agent.memory_stream.id_to_node = {}
                
                # 设置记忆清除标记，防止在退出时保存空记忆
                set_memory_cleared_flag(True)
                
                util.log(1, "已成功清除代理在内存中的记忆")
            
            return True
    except Exception as e:
        util.log(1, f"清除代理记忆时出错: {str(e)}")
        return False

# 用户画像分析锁
portrait_analysis_lock = threading.RLock()
portrait_analysis_time = None

def perform_user_portrait_analysis():
    """
    每晚22点执行的用户画像分析任务
    根据当天对话内容和原有画像，使用LLM生成更新后的用户画像
    """
    global portrait_analysis_time
    global portrait_analysis_lock

    with portrait_analysis_lock:
        if portrait_analysis_time and datetime.datetime.now() - portrait_analysis_time < datetime.timedelta(seconds=60):
            return
        portrait_analysis_time = datetime.datetime.now()

        util.log(1, "开始执行用户画像分析...")

        try:
            # 获取所有用户
            all_users = member_db.new_instance().get_all_users()

            for user in all_users:
                username = user[1]
                try:
                    # 获取当天对话记录
                    today_messages = content_db.new_instance().get_today_messages_by_user(username)

                    # 如果当天没有对话，跳过
                    if not today_messages:
                        util.log(1, f"用户 {username} 今天没有对话记录，跳过画像分析")
                        continue

                    # 构建对话文本
                    conversation_lines = []
                    display_name = "主人" if username == "User" else username
                    for msg_type, msg_content in today_messages:
                        if msg_type in ('member', 'user'):
                            conversation_lines.append(f"{display_name}: {msg_content}")
                        else:
                            conversation_lines.append(f"Fay: {msg_content}")
                    conversation_text = "\n".join(conversation_lines)

                    # 获取原有用户画像
                    current_portrait = member_db.new_instance().get_user_portrait(username)

                    # 构建分析prompt
                    analysis_prompt = f"""你是一个用户画像分析专家。请根据以下信息分析并更新用户画像。

**用户名**: {display_name}

**原有用户画像**:
{current_portrait if current_portrait else "（暂无）"}

**今日对话记录**:
{conversation_text}

**分析要求**:
1. 基于今日对话内容，提取并分析以下维度的信息：
   - 基本信息（姓名、年龄、性别、生日等）
   - 性格特点、兴趣爱好、行为习惯、情感状态
   - 亲朋好友信息（家人、朋友、同事等人物关系）
   - 生活信息（工作、居住、日常活动等）
   - 身体状况（健康状态、疾病、运动习惯等）
   - 身边事物（宠物、车辆、常用物品等）
2. 如果原有画像存在，请在其基础上进行补充和修正
3. 如果发现与原有画像矛盾的信息，以最新对话为准进行更新
4. 画像应简洁明了，使用分点描述，按维度分类整理
5. 只输出用户画像内容，不要输出分析过程
6. 总字数控制在800字以内

请输出更新后的用户画像:"""

                    # 调用LLM进行分析
                    try:
                        response = llm.invoke([
                            SystemMessage(content="你是用户画像分析专家，擅长从对话中提取用户特征。"),
                            HumanMessage(content=analysis_prompt)
                        ])
                        new_portrait = response.content.strip()

                        # 保存新的用户画像
                        member_db.new_instance().update_user_portrait(username, new_portrait)
                        util.log(1, f"用户 {username} 画像分析完成并已保存")

                    except Exception as llm_err:
                        util.log(1, f"用户 {username} LLM分析失败: {llm_err}")

                except Exception as user_err:
                    util.log(1, f"处理用户 {username} 时出错: {user_err}")

            util.log(1, "用户画像分析任务完成")

        except Exception as e:
            util.log(1, f"用户画像分析任务出错: {e}")

# 反思
def perform_daily_reflection():
    global reflection_time
    global reflection_lock
    
    with reflection_lock:
        if reflection_time and datetime.datetime.now() - reflection_time < datetime.timedelta(seconds=60):
            return
        reflection_time = datetime.datetime.now()
 
        # 获取今天的日期，用于确定反思主题
        today = datetime.datetime.now().weekday()
        
        # 根据星期几选择不同反思主题
        reflection_topics = [
            "我与用户的关系发展，以及我如何更好地理解和服务他们",
            "我的知识库如何得到扩展，哪些概念需要进一步理解",
            "我的情感响应模式以及它们如何反映我的核心价值观",
            "我的沟通方式如何影响互动质量，哪些模式最有效",
            "我的行为如何体现我的核心特质，我的自我认知有何变化",
            "今天的经历如何与我的过往记忆建立联系，形成什么样的模式",
            "本周的整体经历与学习"
        ]
        
        # 选择今天的主题(可以按星期轮换或其他逻辑)
        topic = reflection_topics[today % len(reflection_topics)]
        
        # 执行反思，传入当前时间戳
        for username, agent in agents.items():
            try:
                # 获取当前时间作为time_step
                current_time_step = get_current_time_step(username)
                agent.reflect(topic, time_step=current_time_step)
            except KeyError as e:
                util.log(1, f"反思时出现KeyError: {e}，跳过此次反思")
            except Exception as e:
                util.log(1, f"反思时出现错误: {e}，跳过此次反思")
        
        # 记录反思执行情况
        util.log(1, f"反思主题: {topic}")

def save_agent_memory():
    """
    保存代理的记忆到文件
    """
    global agents
    global save_time
    global save_lock
    # 检查记忆清除标记，如果已清除则不保存
    global memory_cleared
    if memory_cleared:
        util.log(1, "检测到记忆已被清除，跳过保存操作")
        return
    
    try:
        # 节流：60s 内只允许保存一次
        with save_lock:
            if save_time and datetime.datetime.now() - save_time < datetime.timedelta(seconds=60):
                return
            save_time = datetime.datetime.now()

        # .memory_cleared 标记文件检查（无需持锁）
        base_dir = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
        mem_base = os.path.join(base_dir, "memory")
        memory_cleared_flag_file = os.path.join(mem_base, ".memory_cleared")
        if os.path.exists(memory_cleared_flag_file):
            util.log(1, "检测到.memory_cleared标记文件，跳过保存操作")
            return

        # 仅短暂持锁拿到用户名快照，避免长时间锁住整个 agents
        with agent_lock:
            usernames = list(agents.keys())

        # 逐个用户保存：每个用户独立持锁，给 question() 等线程留出抢占机会
        for username in usernames:
            try:
                with agent_lock:
                    agent = agents.get(username)
                    if agent is None:
                        continue
                    if agent.memory_stream is None:
                        util.log(1, f"代理 {username} 记忆流未初始化，跳过")
                        continue

                    # 防御性初始化
                    if agent.memory_stream.embeddings is None:
                        agent.memory_stream.embeddings = {}
                    if agent.memory_stream.seq_nodes is None:
                        agent.memory_stream.seq_nodes = []
                    if agent.memory_stream.id_to_node is None:
                        agent.memory_stream.id_to_node = {}
                    if agent.scratch is None:
                        agent.scratch = {}

                    memory_dir = get_user_memory_dir(username)

                    # 完整性检查
                    try:
                        valid_nodes = []
                        for node in agent.memory_stream.seq_nodes:
                            if node is None:
                                continue
                            if not hasattr(node, 'node_id') or not hasattr(node, 'content'):
                                continue
                            raw_content = node.content if isinstance(node.content, str) else str(node.content)
                            cleaned_content = _remove_think_from_text(raw_content)
                            if cleaned_content != raw_content:
                                old_content = raw_content
                                node.content = cleaned_content
                                if (
                                    agent.memory_stream.embeddings is not None
                                    and old_content in agent.memory_stream.embeddings
                                    and cleaned_content not in agent.memory_stream.embeddings
                                ):
                                    agent.memory_stream.embeddings[cleaned_content] = agent.memory_stream.embeddings[old_content]
                            else:
                                node.content = raw_content
                            valid_nodes.append(node)

                        agent.memory_stream.seq_nodes = valid_nodes
                        agent.memory_stream.id_to_node = {node.node_id: node for node in valid_nodes if hasattr(node, 'node_id')}
                        if agent.memory_stream.embeddings is not None:
                            kept_contents = {node.content for node in valid_nodes if hasattr(node, 'content')}
                            agent.memory_stream.embeddings = {
                                key: value
                                for key, value in agent.memory_stream.embeddings.items()
                                if key in kept_contents
                            }
                    except Exception as e:
                        util.log(1, f"检查记忆完整性时出错: {str(e)}")

                    # 保存记忆
                    try:
                        agent.save(memory_dir)
                    except Exception as e:
                        util.log(1, f"调用agent.save()时出错: {str(e)}")
                        try:
                            memory_stream_dir = os.path.join(memory_dir, "memory_stream")
                            os.makedirs(memory_stream_dir, exist_ok=True)

                            with open(os.path.join(memory_stream_dir, "embeddings.json"), "w", encoding='utf-8') as f:
                                json.dump(agent.memory_stream.embeddings or {}, f, ensure_ascii=False, indent=2)

                            with open(os.path.join(memory_stream_dir, "nodes.json"), "w", encoding='utf-8') as f:
                                nodes_data = []
                                for node in agent.memory_stream.seq_nodes:
                                    if node is not None and hasattr(node, 'package'):
                                        try:
                                            nodes_data.append(node.package())
                                        except Exception as node_e:
                                            util.log(1, f"打包节点时出错: {str(node_e)}")
                                json.dump(nodes_data, f, ensure_ascii=False, indent=2)

                            with open(os.path.join(memory_dir, "meta.json"), "w", encoding='utf-8') as f:
                                meta_data = {"id": str(agent.id)} if hasattr(agent, 'id') else {}
                                json.dump(meta_data, f, ensure_ascii=False, indent=2)

                            util.log(1, "通过备用方法成功保存记忆")
                        except Exception as backup_e:
                            util.log(1, f"备用保存方法也失败: {str(backup_e)}")

                    # 更新 scratch
                    try:
                        agent.scratch["first_name"] = cfg.config["attribute"]["name"]
                        agent.scratch["age"] = cfg.config["attribute"]["age"]
                        agent.scratch["sex"] = cfg.config["attribute"]["gender"]
                        agent.scratch["additional"] = cfg.config["attribute"]["additional"]
                        agent.scratch["birthplace"] = cfg.config["attribute"]["birth"]
                        agent.scratch["position"] = cfg.config["attribute"]["position"]
                        agent.scratch["zodiac"] = cfg.config["attribute"]["zodiac"]
                        agent.scratch["constellation"] = cfg.config["attribute"]["constellation"]
                        agent.scratch["contact"] = cfg.config["attribute"]["contact"]
                        agent.scratch["voice"] = cfg.config["attribute"]["voice"]
                        agent.scratch["goal"] = cfg.config["attribute"]["goal"]
                        agent.scratch["occupation"] = cfg.config["attribute"]["job"]
                        agent.scratch["current_time"] = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
                    except Exception as e:
                        util.log(1, f"更新时间时出错: {str(e)}")
            except Exception as e:
                util.log(1, f"保存用户 {username} 的代理记忆失败: {str(e)}")

    except Exception as e:
        util.log(1, f"保存代理记忆失败: {str(e)}")

def get_mcp_tools() -> List[Dict[str, Any]]:
    """Fetch all available MCP tools from the registry."""
    try:
        raw_tools = mcp_runtime.get_enabled_tools() or []
        # 只返回启用的工具，预启动工具只要启用也可以被LLM调用
        filtered = [tool for tool in raw_tools if tool.get("enabled", True)]
        return filtered
    except Exception as e:
        util.log(1, f"Failed to fetch MCP tools: {e}")
        return []
    return []


if __name__ == "__main__":
    init_memory_scheduler()
    for _ in range(3):
        query = "Who is Fay?"
        response = question(query, "User")
        print(f"Q: {query}")
        print(f"A: {response}")
        time.sleep(1)

