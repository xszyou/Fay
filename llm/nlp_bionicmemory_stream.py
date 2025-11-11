# -*- coding: utf-8 -*-
import os
import json
import time
import threading
import requests
import datetime
import schedule
import textwrap
from dataclasses import dataclass
from typing import Any, Callable, Dict, List, Literal, Optional, TypedDict, Tuple
from collections.abc import Mapping, Sequence
from langchain_openai import ChatOpenAI
from langchain_core.messages import HumanMessage, SystemMessage
from langgraph.graph import END, START, StateGraph

# 新增：本地知识库相关导入
import re
from pathlib import Path
import docx
from docx.document import Document
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.table import _Cell, Table
from docx.text.paragraph import Paragraph
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

# 用于处理 .doc 文件的库
try:
    import win32com.client
    WIN32COM_AVAILABLE = True
except ImportError:
    WIN32COM_AVAILABLE = False

from utils import util
import utils.config_util as cfg
from urllib3.exceptions import InsecureRequestWarning
from scheduler.thread_manager import MyThread
from core import content_db
from core import stream_manager
from faymcp import tool_registry as mcp_tool_registry

# 新增：长短期记忆系统相关导入
from bionicmemory.core.chroma_service import ChromaService
from bionicmemory.core.memory_system import LongShortTermMemorySystem, SourceType

os.environ["LANGCHAIN_TRACING_V2"] = "true"
os.environ["LANGCHAIN_ENDPOINT"] = "https://api.smith.langchain.com"
os.environ["LANGCHAIN_API_KEY"] = "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"
os.environ["LANGCHAIN_PROJECT"] = "fay3.11.1_github"

# 加载配置
cfg.load_config()

# 禁用不安全请求警告
requests.packages.urllib3.disable_warnings(category=InsecureRequestWarning)

# 记忆系统全局变量
chroma_service = None  # ChromaDB服务实例
memory_system = None   # 长短期记忆系统实例
memory_system_lock = threading.RLock()  # 保护记忆系统的锁

# 当前会话用户名（保留，用于兼容性）
current_username = None

llm = ChatOpenAI(
        model=cfg.gpt_model_engine,
        base_url=cfg.gpt_base_url,
        api_key=cfg.key_gpt_api_key,
        streaming=True
    )


def init_memory_system():
    """
    初始化长短期记忆系统

    Returns:
        bool: 是否初始化成功
    """
    global chroma_service, memory_system

    try:
        util.log(1, "正在初始化记忆系统...")

        # 初始化ChromaDB服务
        chroma_service = ChromaService()
        if not chroma_service:
            util.log(1, "ChromaDB服务初始化失败")
            return False

        # 初始化长短期记忆系统
        memory_system = LongShortTermMemorySystem(
            chroma_service=chroma_service,
            summary_threshold=500,
            max_retrieval_results=10,
            cluster_multiplier=3,
            retrieval_multiplier=2
        )

        util.log(1, "记忆系统初始化成功")
        return True

    except Exception as e:
        util.log(1, f"记忆系统初始化失败: {e}")
        return False


# 在模块加载时初始化记忆系统
init_memory_system()


@dataclass
class WorkflowToolSpec:
    name: str
    description: str
    schema: Dict[str, Any]
    executor: Callable[[Dict[str, Any], int], Tuple[bool, Optional[str], Optional[str]]]
    example_args: Dict[str, Any]


class ToolCall(TypedDict):
    name: str
    args: Dict[str, Any]


class ToolResult(TypedDict, total=False):
    call: ToolCall
    success: bool
    output: Optional[str]
    error: Optional[str]
    attempt: int


class ConversationMessage(TypedDict):
    role: Literal["user", "assistant"]
    content: str


class AgentState(TypedDict, total=False):
    request: str
    messages: List[ConversationMessage]
    tool_results: List[ToolResult]
    next_action: Optional[ToolCall]
    status: Literal["planning", "needs_tool", "completed", "failed"]
    final_response: Optional[str]
    final_messages: Optional[List[SystemMessage | HumanMessage]]
    planner_preview: Optional[str]
    audit_log: List[str]
    context: Dict[str, Any]
    error: Optional[str]
    max_steps: int


def _truncate_text(text: Any, limit: int = 400) -> str:
    text_str = "" if text is None else str(text)
    if len(text_str) <= limit:
        return text_str
    return text_str[:limit] + "..."


def _extract_text_from_result(value: Any, *, depth: int = 0) -> List[str]:
    """Try to pull human-readable text snippets from tool results."""
    if value is None:
        return []
    if depth > 6:
        return []
    if isinstance(value, (str, int, float, bool)):
        text = str(value).strip()
        return [text] if text else []
    if isinstance(value, Mapping):
        # Prefer explicit text/content fields
        if "text" in value and not isinstance(value["text"], (dict, list, tuple)):
            text = str(value["text"]).strip()
            return [text] if text else []
        if "content" in value:
            segments: List[str] = []
            for item in value.get("content", []):
                segments.extend(_extract_text_from_result(item, depth=depth + 1))
            if segments:
                return segments
        segments = []
        for key, item in value.items():
            if key in {"meta", "annotations", "uid", "id", "messageId"}:
                continue
            segments.extend(_extract_text_from_result(item, depth=depth + 1))
        return segments
    if isinstance(value, Sequence) and not isinstance(value, (bytes, bytearray)):
        segments: List[str] = []
        for item in value:
            segments.extend(_extract_text_from_result(item, depth=depth + 1))
        return segments
    if hasattr(value, "text") and not callable(getattr(value, "text")):
        text = str(getattr(value, "text", "")).strip()
        return [text] if text else []
    if hasattr(value, "__dict__"):
        return _extract_text_from_result(vars(value), depth=depth + 1)
    text = str(value).strip()
    return [text] if text else []


def _normalize_tool_output(result: Any) -> str:
    """Convert structured tool output to a concise human-readable string."""
    if result is None:
        return ""
    segments = _extract_text_from_result(result)
    if segments:
        cleaned = [segment for segment in segments if segment]
        if cleaned:
            return "\n".join(dict.fromkeys(cleaned))
    try:
        return json.dumps(result, ensure_ascii=False, default=lambda o: getattr(o, "__dict__", str(o)))
    except TypeError:
        return str(result)


def _truncate_history(history: List[ToolResult], limit: int = 6) -> str:
    if not history:
        return "（暂无）"
    lines: List[str] = []
    for item in history[-limit:]:
        call = item.get("call", {})
        name = call.get("name", "未知工具")
        attempt = item.get("attempt", 0)
        success = item.get("success", False)
        status = "成功" if success else "失败"
        lines.append(f"- {name} 第 {attempt} 次 → {status}")
        if item.get("output"):
            lines.append("  输出：" + _truncate_text(item["output"], 200))
        if item.get("error"):
            lines.append("  错误：" + _truncate_text(item["error"], 200))
    return "\n".join(lines)


def _format_schema_parameters(schema: Dict[str, Any]) -> List[str]:
    if not schema:
        return ["  - 无参数"]
    props = schema.get("properties") or {}
    if not props:
        return ["  - 无参数"]
    required = set(schema.get("required") or [])
    lines: List[str] = []
    for field, meta in props.items():
        meta = meta or {}
        field_type = meta.get("type", "string")
        desc = (meta.get("description") or "").strip()
        req_label = "必填" if field in required else "可选"
        line = f"  - {field} ({field_type}，{req_label})"
        if desc:
            line += f"：{desc}"
        lines.append(line)
    return lines or ["  - 无参数"]


def _generate_example_args(schema: Dict[str, Any]) -> Dict[str, Any]:
    example: Dict[str, Any] = {}
    if not schema:
        return example
    props = schema.get("properties") or {}
    for field, meta in props.items():
        meta = meta or {}
        if "default" in meta:
            example[field] = meta["default"]
            continue
        enum_values = meta.get("enum") or []
        if enum_values:
            example[field] = enum_values[0]
            continue
        field_type = meta.get("type", "string")
        if field_type in ("number", "integer"):
            example[field] = 0
        elif field_type == "boolean":
            example[field] = True
        elif field_type == "array":
            example[field] = []
        elif field_type == "object":
            example[field] = {}
        else:
            description_hint = meta.get("description") or ""
            example[field] = description_hint or ""
    return example


def _format_tool_block(spec: WorkflowToolSpec) -> str:
    param_lines = _format_schema_parameters(spec.schema)
    example = json.dumps(spec.example_args, ensure_ascii=False) if spec.example_args else "{}"
    lines = [
        f"- 工具名：{spec.name}",
        f"  功能：{spec.description or '暂无描述'}",
        "  参数：",
        *param_lines,
        f"  示例：{example}",
    ]
    return "\n".join(lines)


def _build_workflow_tool_spec(tool_def: Dict[str, Any]) -> Optional[WorkflowToolSpec]:
    if not tool_def:
        return None
    name = tool_def.get("name")
    if not name:
        return None
    description = tool_def.get("description") or tool_def.get("summary") or ""
    schema = tool_def.get("inputSchema") or {}
    example_args = _generate_example_args(schema)

    def _executor(args: Dict[str, Any], attempt: int) -> Tuple[bool, Optional[str], Optional[str]]:
        try:
            resp = requests.post(
                f"http://127.0.0.1:5010/api/mcp/tools/{name}",
                json=args,
                timeout=120,
            )
            resp.raise_for_status()
            data = resp.json()
        except Exception as exc:
            util.log(1, f"调用工具 {name} 异常: {exc}")
            return False, None, str(exc)

        if data.get("success"):
            result = data.get("result")
            output = _normalize_tool_output(result)
            return True, output, None

        error_msg = data.get("error") or "未知错误"
        util.log(1, f"调用工具 {name} 失败: {error_msg}")
        return False, None, error_msg

    return WorkflowToolSpec(
        name=name,
        description=description,
        schema=schema,
        executor=_executor,
        example_args=example_args,
    )


def _format_tools_for_prompt(tool_specs: Dict[str, WorkflowToolSpec]) -> str:
    if not tool_specs:
        return "（暂无可用工具）"
    return "\n".join(_format_tool_block(spec) for spec in tool_specs.values())


def _build_planner_messages(state: AgentState) -> List[SystemMessage | HumanMessage]:
    context = state.get("context", {}) or {}
    system_prompt = context.get("system_prompt", "")
    request = state.get("request", "")
    tool_specs = context.get("tool_registry", {}) or {}
    planner_preview = state.get("planner_preview")
    conversation = state.get("messages", []) or []
    history = state.get("tool_results", []) or []
    knowledge_context = context.get("knowledge_context", "")
    observation = context.get("observation", "")

    convo_text = "\n".join(f"{msg['role']}: {msg['content']}" for msg in conversation) or "（暂无对话）"
    history_text = _truncate_history(history)
    tools_text = _format_tools_for_prompt(tool_specs)
    preview_section = f"\n（规划器预览：{planner_preview}）" if planner_preview else ""

    user_block = textwrap.dedent(
        f"""

**当前请求**
{request}

{system_prompt}

**额外观察**
{observation or '（无补充）'}

**相关知识**
{knowledge_context or '（无相关知识）'}

**可用工具**
{tools_text}

**历史工具执行**
{history_text}{preview_section}

请返回 JSON，格式如下：
- 若需要调用工具：
    {{"action": "tool", "tool": "工具名", "args": {{...}}}}
- 若直接回复：
    {{"action": "finish_text"}}

对话及工具记录：
{convo_text}
        """
    ).strip()

    return [
        SystemMessage(content="你负责规划下一步行动，请严格输出合法 JSON。"),
        HumanMessage(content=user_block),
    ]


def _build_final_messages(state: AgentState) -> List[SystemMessage | HumanMessage]:
    context = state.get("context", {}) or {}
    system_prompt = context.get("system_prompt", "")
    request = state.get("request", "")
    knowledge_context = context.get("knowledge_context", "")
    observation = context.get("observation", "")
    conversation = state.get("messages", []) or []
    planner_preview = state.get("planner_preview")
    conversation_block = "\n".join(f"{msg['role']}: {msg['content']}" for msg in conversation) or "（暂无对话）"
    history_text = _truncate_history(state.get("tool_results", []))
    preview_section = f"\n（规划器建议：{planner_preview}）" if planner_preview else ""

    user_block = textwrap.dedent(
        f"""
**当前请求**
{request}

{system_prompt}

**相关知识**
{knowledge_context or '（无相关知识）'}

**其他观察**
{observation or '（无补充）'}

**工具执行摘要**
{history_text}{preview_section}

**对话及工具记录**
{conversation_block}
        """
    ).strip()

    return [
        SystemMessage(content="你是最终回复的口播助手，请用中文自然表达。"),
        HumanMessage(content=user_block),
    ]


def _call_planner_llm(state: AgentState) -> Dict[str, Any]:
    response = llm.invoke(_build_planner_messages(state))
    content = getattr(response, "content", None)
    if not isinstance(content, str):
        raise RuntimeError("规划器返回内容异常，未获得字符串。")
    trimmed = content.strip()
    try:
        decision = json.loads(trimmed)
    except json.JSONDecodeError as exc:
        raise RuntimeError(f"规划器返回的 JSON 无法解析: {trimmed}") from exc
    decision.setdefault("_raw", trimmed)
    return decision


def _plan_next_action(state: AgentState) -> AgentState:
    context = state.get("context", {}) or {}
    audit_log = list(state.get("audit_log", []))
    history = state.get("tool_results", []) or []
    max_steps = state.get("max_steps", 12)
    if len(history) >= max_steps:
        audit_log.append("规划器：超过最大步数，终止流程。")
        return {
            "status": "failed",
            "audit_log": audit_log,
            "error": "工具调用步数超限",
            "context": context,
        }

    decision = _call_planner_llm(state)
    audit_log.append(f"规划器：决策 -> {decision.get('_raw', decision)}")

    action = decision.get("action")
    if action == "tool":
        tool_name = decision.get("tool")
        tool_registry: Dict[str, WorkflowToolSpec] = context.get("tool_registry", {})
        if tool_name not in tool_registry:
            audit_log.append(f"规划器：未知工具 {tool_name}")
            return {
                "status": "failed",
                "audit_log": audit_log,
                "error": f"未知工具 {tool_name}",
                "context": context,
            }
        args = decision.get("args") or {}

        if history:
            last_entry = history[-1]
            last_call = last_entry.get("call", {}) or {}
            if (
                last_entry.get("success")
                and last_call.get("name") == tool_name
                and (last_call.get("args") or {}) == args
                and last_entry.get("output")
            ):
                recent_attempts = sum(
                    1
                    for item in reversed(history)
                    if item.get("call", {}).get("name") == tool_name
                )
                if recent_attempts >= 1:
                    audit_log.append(
                        "规划器：检测到工具重复调用，使用最新结果产出最终回复。"
                    )
                    final_messages = _build_final_messages(state)
                    preview = last_entry.get("output")
                    return {
                        "status": "completed",
                        "planner_preview": preview,
                        "final_response": None,
                        "final_messages": final_messages,
                        "audit_log": audit_log,
                        "context": context,
                    }
        return {
            "next_action": {"name": tool_name, "args": args},
            "status": "needs_tool",
            "audit_log": audit_log,
            "context": context,
        }

    if action in {"finish", "finish_text"}:
        preview = decision.get("message")
        final_messages = _build_final_messages(state)
        audit_log.append("规划器：任务完成，准备输出最终回复。")
        return {
            "status": "completed",
            "planner_preview": preview,
            "final_response": preview if action == "finish" else None,
            "final_messages": final_messages,
            "audit_log": audit_log,
            "context": context,
        }

    raise RuntimeError(f"未知的规划器决策: {decision}")


def _execute_tool(state: AgentState) -> AgentState:
    context = dict(state.get("context", {}) or {})
    action = state.get("next_action")
    if not action:
        return {
            "status": "failed",
            "error": "缺少要执行的工具指令",
            "context": context,
        }

    history = list(state.get("tool_results", []) or [])
    audit_log = list(state.get("audit_log", []) or [])
    conversation = list(state.get("messages", []) or [])

    name = action.get("name")
    args = action.get("args", {})
    tool_registry: Dict[str, WorkflowToolSpec] = context.get("tool_registry", {})
    spec = tool_registry.get(name)
    if not spec:
        return {
            "status": "failed",
            "error": f"未知工具 {name}",
            "context": context,
        }

    attempts = sum(1 for item in history if item.get("call", {}).get("name") == name)
    success, output, error = spec.executor(args, attempts)
    result: ToolResult = {
        "call": {"name": name, "args": args},
        "success": success,
        "output": output,
        "error": error,
        "attempt": attempts + 1,
    }
    history.append(result)
    audit_log.append(f"执行器：{name} 第 {result['attempt']} 次 -> {'成功' if success else '失败'}")

    message_lines = [
        f"[TOOL] {name} {'成功' if success else '失败'}。",
    ]
    if output:
        message_lines.append(f"[TOOL] 输出：{_truncate_text(output, 200)}")
    if error:
        message_lines.append(f"[TOOL] 错误：{_truncate_text(error, 200)}")
    conversation.append({"role": "assistant", "content": "\n".join(message_lines)})

    return {
        "tool_results": history,
        "messages": conversation,
        "next_action": None,
        "audit_log": audit_log,
        "status": "planning",
        "error": error if not success else None,
        "context": context,
    }


def _route_decision(state: AgentState) -> str:
    return "call_tool" if state.get("status") == "needs_tool" else "end"


def _build_workflow_app() -> StateGraph:
    graph = StateGraph(AgentState)
    graph.add_node("plan_next", _plan_next_action)
    graph.add_node("call_tool", _execute_tool)
    graph.add_edge(START, "plan_next")
    graph.add_conditional_edges(
        "plan_next",
        _route_decision,
        {
            "call_tool": "call_tool",
            "end": END,
        },
    )
    graph.add_edge("call_tool", "plan_next")
    return graph.compile()


_WORKFLOW_APP = _build_workflow_app()

# 新增：本地知识库相关函数
def read_doc_file(file_path):
    """
    读取doc文件内容
    
    参数:
        file_path: doc文件路径
        
    返回:
        str: 文档内容
    """
    try:
        # 方法1: 使用 win32com.client（Windows系统，推荐用于.doc文件）
        if WIN32COM_AVAILABLE:
            word = None
            doc = None
            try:
                import pythoncom
                pythoncom.CoInitialize()  # 初始化COM组件
                
                word = win32com.client.Dispatch("Word.Application")
                word.Visible = False
                doc = word.Documents.Open(file_path)
                content = doc.Content.Text
                
                # 先保存内容，再尝试关闭
                if content and content.strip():
                    try:
                        doc.Close()
                        word.Quit()
                    except Exception as close_e:
                        util.log(1, f"关闭Word应用程序时出错: {str(close_e)}，但内容已成功提取")
                    
                    try:
                        pythoncom.CoUninitialize()  # 清理COM组件
                    except:
                        pass
                    
                    return content.strip()
                
            except Exception as e:
                util.log(1, f"使用 win32com 读取 .doc 文件失败: {str(e)}")
            finally:
                # 确保资源被释放
                try:
                    if doc:
                        doc.Close()
                except:
                    pass
                try:
                    if word:
                        word.Quit()
                except:
                    pass
                try:
                    pythoncom.CoUninitialize()
                except:
                    pass
        
        # 方法2: 简单的二进制文本提取（备选方案）
        try:
            with open(file_path, 'rb') as f:
                raw_data = f.read()
                # 尝试提取可打印的文本
                text_parts = []
                current_text = ""
                
                for byte in raw_data:
                    char = chr(byte) if 32 <= byte <= 126 or byte in [9, 10, 13] else None
                    if char:
                        current_text += char
                    else:
                        if len(current_text) > 3:  # 只保留长度大于3的文本片段
                            text_parts.append(current_text.strip())
                        current_text = ""
                
                if len(current_text) > 3:
                    text_parts.append(current_text.strip())
                
                # 过滤和清理文本
                filtered_parts = []
                for part in text_parts:
                    # 移除过多的重复字符和无意义的片段
                    if (len(part) > 5 and 
                        not part.startswith('Microsoft') and 
                        not all(c in '0123456789-_.' for c in part) and
                        len(set(part)) > 3):  # 字符种类要多样
                        filtered_parts.append(part)
                
                if filtered_parts:
                    return '\n'.join(filtered_parts)
                    
        except Exception as e:
            util.log(1, f"使用二进制方法读取 .doc 文件失败: {str(e)}")
        
        util.log(1, f"无法读取 .doc 文件 {file_path}，建议转换为 .docx 格式")
        return ""
        
    except Exception as e:
        util.log(1, f"读取doc文件 {file_path} 时出错: {str(e)}")
        return ""

def read_docx_file(file_path):
    """
    读取docx文件内容
    
    参数:
        file_path: docx文件路径
        
    返回:
        str: 文档内容
    """
    try:
        doc = docx.Document(file_path)
        content = []
        
        for element in doc.element.body:
            if isinstance(element, CT_P):
                paragraph = Paragraph(element, doc)
                if paragraph.text.strip():
                    content.append(paragraph.text.strip())
            elif isinstance(element, CT_Tbl):
                table = Table(element, doc)
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        content.append(" | ".join(row_text))
        
        return "\n".join(content)
    except Exception as e:
        util.log(1, f"读取docx文件 {file_path} 时出错: {str(e)}")
        return ""
    
def read_pptx_file(file_path):
    """
    读取pptx文件内容
    
    参数:
        file_path: pptx文件路径
        
    返回:
        str: 演示文稿内容
    """
    if not PPTX_AVAILABLE:
        util.log(1, "python-pptx 库未安装，无法读取 PowerPoint 文件")
        return ""
        
    try:
        prs = Presentation(file_path)
        content = []
        
        for i, slide in enumerate(prs.slides):
            slide_content = [f"第{i+1}页："]
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text.strip())
                    
            if len(slide_content) > 1:  # 有内容才添加
                content.append("\n".join(slide_content))
        
        return "\n\n".join(content)
    except Exception as e:
        util.log(1, f"读取pptx文件 {file_path} 时出错: {str(e)}")
        return ""

def load_local_knowledge_base():
    """
    加载本地知识库内容
    
    返回:
        dict: 文件名到内容的映射
    """
    knowledge_base = {}
    
    # 获取llm/data目录路径
    current_dir = os.path.dirname(os.path.abspath(__file__))
    data_dir = os.path.join(current_dir, "data")
    
    if not os.path.exists(data_dir):
        util.log(1, f"知识库目录不存在: {data_dir}")
        return knowledge_base
    
    # 遍历data目录中的文件
    for file_path in Path(data_dir).iterdir():
        if not file_path.is_file():
            continue
            
        file_name = file_path.name
        file_extension = file_path.suffix.lower()
        
        try:
            if file_extension == '.docx':
                content = read_docx_file(str(file_path))
            elif file_extension == '.doc':
                content = read_doc_file(str(file_path))
            elif file_extension == '.pptx':
                content = read_pptx_file(str(file_path))
            else:
                # 尝试作为文本文件读取
                try:
                    with open(file_path, 'r', encoding='utf-8') as f:
                        content = f.read()
                except UnicodeDecodeError:
                    try:
                        with open(file_path, 'r', encoding='gbk') as f:
                            content = f.read()
                    except UnicodeDecodeError:
                        util.log(1, f"无法解码文件: {file_name}")
                        continue
            
            if content.strip():
                knowledge_base[file_name] = content
                util.log(1, f"成功加载知识库文件: {file_name} ({len(content)} 字符)")
            
        except Exception as e:
            util.log(1, f"加载知识库文件 {file_name} 时出错: {str(e)}")
    
    return knowledge_base

def search_knowledge_base(query, knowledge_base, max_results=3):
    """
    在知识库中搜索相关内容
    
    参数:
        query: 查询内容
        knowledge_base: 知识库字典
        max_results: 最大返回结果数
        
    返回:
        list: 相关内容列表
    """
    if not knowledge_base:
        return []
    
    results = []
    query_lower = query.lower()
    
    # 搜索关键词
    query_keywords = re.findall(r'\w+', query_lower)
    
    for file_name, content in knowledge_base.items():
        content_lower = content.lower()
        
        # 计算匹配度
        score = 0
        matched_sentences = []
        
        # 按句子分割内容
        sentences = re.split(r'[。！？\n]', content)
        
        for sentence in sentences:
            if not sentence.strip():
                continue
                
            sentence_lower = sentence.lower()
            sentence_score = 0
            
            # 计算关键词匹配度
            for keyword in query_keywords:
                if keyword in sentence_lower:
                    sentence_score += 1
            
            # 如果句子有匹配，记录
            if sentence_score > 0:
                matched_sentences.append((sentence.strip(), sentence_score))
                score += sentence_score
        
        # 如果有匹配的内容
        if score > 0:
            # 按匹配度排序句子
            matched_sentences.sort(key=lambda x: x[1], reverse=True)
            
            # 取前几个最相关的句子
            relevant_sentences = [sent[0] for sent in matched_sentences[:5] if sent[0]]
            
            if relevant_sentences:
                results.append({
                    'file_name': file_name,
                    'score': score,
                    'content': '\n'.join(relevant_sentences)
                })
    
    # 按匹配度排序
    results.sort(key=lambda x: x['score'], reverse=True)
    
    return results[:max_results]

# 全局知识库缓存
_knowledge_base_cache = None
_knowledge_base_load_time = None
_knowledge_base_file_times = {}  # 存储文件的最后修改时间

def check_knowledge_base_changes():
    """
    检查知识库文件是否有变化
    
    返回:
        bool: 如果有文件变化返回True，否则返回False
    """
    global _knowledge_base_file_times
    
    # 获取llm/data目录路径
    current_dir = os.path.dirname(os.path.abspath(__file__))
    data_dir = os.path.join(current_dir, "data")
    
    if not os.path.exists(data_dir):
        return False
    
    current_file_times = {}
    
    # 遍历data目录中的文件
    for file_path in Path(data_dir).iterdir():
        if not file_path.is_file():
            continue
        
        file_name = file_path.name
        file_extension = file_path.suffix.lower()
        
        # 只检查支持的文件格式
        if file_extension in ['.docx', '.doc', '.pptx', '.txt'] or file_extension == '':
            try:
                mtime = os.path.getmtime(str(file_path))
                current_file_times[file_name] = mtime
            except OSError:
                continue
    
    # 检查是否有变化
    if not _knowledge_base_file_times:
        # 第一次检查，保存文件时间
        _knowledge_base_file_times = current_file_times
        return True
    
    # 比较文件时间
    if set(current_file_times.keys()) != set(_knowledge_base_file_times.keys()):
        # 文件数量发生变化
        _knowledge_base_file_times = current_file_times
        return True
    
    for file_name, mtime in current_file_times.items():
        if file_name not in _knowledge_base_file_times or _knowledge_base_file_times[file_name] != mtime:
            # 文件被修改
            _knowledge_base_file_times = current_file_times
            return True
    
    return False

def init_knowledge_base():
    """
    初始化知识库，在系统启动时调用
    """
    global _knowledge_base_cache, _knowledge_base_load_time
    
    util.log(1, "初始化本地知识库...")
    _knowledge_base_cache = load_local_knowledge_base()
    _knowledge_base_load_time = time.time()
    
    # 初始化文件修改时间跟踪
    check_knowledge_base_changes()
    
    util.log(1, f"知识库初始化完成，共 {len(_knowledge_base_cache)} 个文件")

def get_knowledge_base():
    """
    获取知识库，使用缓存机制
    
    返回:
        dict: 知识库内容
    """
    global _knowledge_base_cache, _knowledge_base_load_time
    
    # 如果缓存为空，先初始化
    if _knowledge_base_cache is None:
        init_knowledge_base()
        return _knowledge_base_cache
    
    # 检查文件是否有变化
    if check_knowledge_base_changes():
        util.log(1, "检测到知识库文件变化，正在重新加载...")
        _knowledge_base_cache = load_local_knowledge_base()
        _knowledge_base_load_time = time.time()
        util.log(1, f"知识库重新加载完成，共 {len(_knowledge_base_cache)} 个文件")
    
    return _knowledge_base_cache


def question(content, username, observation=None):
    """处理用户提问并返回回复。"""
    global current_username
    current_username = username
    full_response_text = ""
    accumulated_text = ""
    default_punctuations = [",", ".", "!", "?", "\n", "\uFF0C", "\u3002", "\uFF01", "\uFF1F"]
    is_first_sentence = True

    from core import stream_manager
    sm = stream_manager.new_instance()
    conversation_id = sm.get_conversation_id(username)

    # 记忆系统已在全局初始化，无需创建agent
    # 直接从配置文件获取人物设定
    agent_desc = {
        "first_name": cfg.config["attribute"]["name"],
        "last_name": "",
        "age": cfg.config["attribute"]["age"],
        "sex": cfg.config["attribute"]["gender"],
        "additional": cfg.config["attribute"]["additional"],
        "birthplace": cfg.config["attribute"]["birth"],
        "position": cfg.config["attribute"]["position"],
        "zodiac": cfg.config["attribute"]["zodiac"],
        "constellation": cfg.config["attribute"]["constellation"],
        "contact": cfg.config["attribute"]["contact"],
        "voice": cfg.config["attribute"]["voice"],
        "goal": cfg.config["attribute"]["goal"],
        "occupation": cfg.config["attribute"]["job"],
        "current_time": datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    }

    # 使用新记忆系统处理用户消息
    # 一次性完成：入库 → 长期检索 → 短期检索 → 生成提示语
    short_term_records = []
    memory_prompt = ""
    query_embedding = None

    try:
        short_term_records, memory_prompt, query_embedding = memory_system.process_user_message(
            content, user_id=username
        )
        util.log(1, f"记忆检索成功，获取 {len(short_term_records)} 条相关记录")
    except Exception as exc:
        util.log(1, f"记忆检索失败: {exc}")
        # 失败时使用空值，不影响后续流程
        short_term_records = []
        memory_prompt = ""
        query_embedding = None

    knowledge_context = ""
    try:
        knowledge_base = get_knowledge_base()
        if knowledge_base:
            knowledge_results = search_knowledge_base(content, knowledge_base, max_results=3)
            if knowledge_results:
                parts = ["**本地知识库相关信息**："]
                for result in knowledge_results:
                    parts.append(f"来源文件：{result['file_name']}")
                    parts.append(result["content"])
                    parts.append("")
                knowledge_context = "\n".join(parts).strip()
                util.log(1, f"找到 {len(knowledge_results)} 条相关知识库信息")
    except Exception as exc:
        util.log(1, f"搜索知识库时出错: {exc}")

    # 方案B：保留人设信息，补充记忆提示语
    # 1. 构建人设部分
    persona_prompt = f"""\n**角色设定**\n
- 名字：{agent_desc['first_name']}
- 性别：{agent_desc['sex']}
- 年龄：{agent_desc['age']}
- 职业：{agent_desc['occupation']}
- 出生地：{agent_desc['birthplace']}
- 星座：{agent_desc['constellation']}
- 生肖：{agent_desc['zodiac']}
- 联系方式：{agent_desc['contact']}
- 定位：{agent_desc['position']}
- 目标：{agent_desc['goal']}
- 补充信息：{agent_desc['additional']}\n

"""

    # 2. 合并人设和记忆提示语
    if memory_prompt:
        system_prompt =  memory_prompt + persona_prompt
    else:
        # 如果记忆系统返回空提示语，使用基础提示语
        system_prompt = persona_prompt + "请根据用户的问题，提供有帮助的回答。"

    try:
        history_records = content_db.new_instance().get_recent_messages_by_user(username=username, limit=30)
    except Exception as exc:
        util.log(1, f"加载历史消息失败: {exc}")
        history_records = []

    messages_buffer: List[ConversationMessage] = []

    def append_to_buffer(role: str, text_value: str) -> None:
        if not text_value:
            return
        messages_buffer.append({"role": role, "content": text_value})
        if len(messages_buffer) > 60:
            del messages_buffer[:-60]

    for msg_type, msg_text in history_records:
        role = 'assistant'
        if msg_type and msg_type.lower() in ('member', 'user'):
            role = 'user'
        append_to_buffer(role, msg_text)

    if (
        not messages_buffer
        or messages_buffer[-1]['role'] != 'user'
        or messages_buffer[-1]['content'] != content
    ):
        append_to_buffer('user', content)

    tool_registry: Dict[str, WorkflowToolSpec] = {}
    try:
        mcp_tools = get_mcp_tools()
    except Exception as exc:
        util.log(1, f"获取工具列表失败: {exc}")
        mcp_tools = []
    for tool_def in mcp_tools:
        spec = _build_workflow_tool_spec(tool_def)
        if spec:
            tool_registry[spec.name] = spec

    try:
        from utils.stream_state_manager import get_state_manager as _get_state_manager

        state_mgr = _get_state_manager()
        session_label = "workflow_agent" if tool_registry else "llm_stream"
        if not state_mgr.is_session_active(username, conversation_id=conversation_id):
            state_mgr.start_new_session(username, session_label, conversation_id=conversation_id)
    except Exception:
        state_mgr = None

    try:
        from utils.stream_text_processor import get_processor

        processor = get_processor()
        punctuation_list = getattr(processor, "punctuation_marks", default_punctuations)
    except Exception:
        processor = None
        punctuation_list = default_punctuations
    def write_sentence(text: str, *, force_first: bool = False, force_end: bool = False) -> None:
        if text is None:
            text = ""
        if not isinstance(text, str):
            text = str(text)
        if not text and not force_end and not force_first:
            return
        marked_text = None
        if state_mgr is not None:
            try:
                marked_text, _, _ = state_mgr.prepare_sentence(
                    username,
                    text,
                    force_first=force_first,
                    force_end=force_end,
                    conversation_id=conversation_id,
                )
            except Exception:
                marked_text = None
        if marked_text is None:
            prefix = "_<isfirst>" if force_first else ""
            suffix = "_<isend>" if force_end else ""
            marked_text = f"{prefix}{text}{suffix}"
        stream_manager.new_instance().write_sentence(username, marked_text, conversation_id=conversation_id)

    def stream_response_chunks(chunks, prepend_text: str = "") -> None:
        nonlocal accumulated_text, full_response_text, is_first_sentence
        if prepend_text:
            accumulated_text += prepend_text
            full_response_text += prepend_text
        for chunk in chunks:
            if sm.should_stop_generation(username, conversation_id=conversation_id):
                util.log(1, f"检测到停止标志，中断文本生成: {username}")
                break
            if isinstance(chunk, str):
                flush_text = chunk
            elif isinstance(chunk, dict):
                flush_text = chunk.get("content")
            else:
                flush_text = getattr(chunk, "content", None)
            if isinstance(flush_text, list):
                flush_text = "".join(part if isinstance(part, str) else "" for part in flush_text)
            if not flush_text:
                continue
            flush_text = str(flush_text)
            accumulated_text += flush_text
            full_response_text += flush_text
            if len(accumulated_text) >= 20:
                while True:
                    last_punct_pos = -1
                    for punct in punctuation_list:
                        pos = accumulated_text.rfind(punct)
                        if pos > last_punct_pos:
                            last_punct_pos = pos
                    if last_punct_pos > 10:
                        sentence_text = accumulated_text[: last_punct_pos + 1]
                        write_sentence(sentence_text, force_first=is_first_sentence)
                        is_first_sentence = False
                        accumulated_text = accumulated_text[last_punct_pos + 1 :].lstrip()
                    else:
                        break

    def finalize_stream(force_end: bool = False) -> None:
        nonlocal accumulated_text, is_first_sentence
        if accumulated_text:
            write_sentence(accumulated_text, force_first=is_first_sentence, force_end=force_end)
            is_first_sentence = False
            accumulated_text = ""
        elif force_end:
            if state_mgr is not None:
                try:
                    session_info = state_mgr.get_session_info(username, conversation_id=conversation_id)
                except Exception:
                    session_info = None
                if not session_info or not session_info.get("is_end_sent", False):
                    write_sentence("", force_end=True)
            else:
                write_sentence("", force_end=True)

    def run_workflow(tool_registry: Dict[str, WorkflowToolSpec]) -> bool:
        nonlocal accumulated_text, full_response_text, is_first_sentence, messages_buffer

        initial_state: AgentState = {
            "request": content,
            "messages": messages_buffer,
            "tool_results": [],
            "audit_log": [],
            "status": "planning",
            "max_steps": 30,
            "context": {
                "system_prompt": system_prompt,
                "knowledge_context": knowledge_context,
                "observation": observation,
                "tool_registry": tool_registry,
            },
        }
        
        config = {"configurable": {"thread_id": f"workflow-{username}-{conversation_id}"}}
        workflow_app = _WORKFLOW_APP
        is_agent_think_start = False
        final_state: Optional[AgentState] = None
        final_stream_done = False

        try:
            for event in workflow_app.stream(initial_state, config=config, stream_mode="updates"):
                if sm.should_stop_generation(username, conversation_id=conversation_id):
                    util.log(1, f"检测到停止标志，中断工作流生成: {username}")
                    break
                step, state = next(iter(event.items()))
                final_state = state
                status = state.get("status")

                state_messages = state.get("messages") or []
                if state_messages and len(state_messages) > len(messages_buffer):
                    messages_buffer.extend(state_messages[len(messages_buffer):])
                    if len(messages_buffer) > 60:
                        del messages_buffer[:-60]

                if step == "plan_next":
                    if status == "needs_tool":
                        next_action = state.get("next_action") or {}
                        tool_name = next_action.get("name") or "unknown_tool"
                        tool_args = next_action.get("args") or {}
                        audit_log = state.get("audit_log") or []
                        decision_note = audit_log[-1] if audit_log else ""
                        if "->" in decision_note:
                            decision_note = decision_note.split("->", 1)[1].strip()
                        args_text = json.dumps(tool_args, ensure_ascii=False)
                        message_lines = [
                            "[PLAN] Planner preparing to call a tool.",
                            f"[PLAN] Decision: {decision_note}" if decision_note else "[PLAN] Decision: (missing)",
                            f"[PLAN] Tool: {tool_name}",
                            f"[PLAN] Args: {args_text}",
                        ]
                        message = "\n".join(message_lines) + "\n"
                        if not is_agent_think_start:
                            message = "<think>" + message
                            is_agent_think_start = True
                        write_sentence(message, force_first=is_first_sentence)
                        is_first_sentence = False
                        full_response_text += message
                        append_to_buffer('assistant', message.strip())
                    elif status == "completed" and not final_stream_done:
                        closing = "</think>" if is_agent_think_start else ""
                        final_messages = state.get("final_messages")
                        final_response = state.get("final_response")
                        success = False
                        if final_messages:
                            try:
                                stream_response_chunks(llm.stream(final_messages), prepend_text=closing)
                                success = True
                            except requests.exceptions.RequestException as stream_exc:
                                util.log(1, f"最终回复流式输出失败: {stream_exc}")
                        elif final_response:
                            stream_response_chunks([closing + final_response])
                            success = True
                        elif closing:
                            accumulated_text += closing
                            full_response_text += closing
                        final_stream_done = success
                        is_agent_think_start = False
                elif step == "call_tool":
                    history = state.get("tool_results") or []
                    if history:
                        last = history[-1]
                        call_info = last.get("call", {}) or {}
                        tool_name = call_info.get("name") or "unknown_tool"
                        success = last.get("success", False)
                        status_text = "SUCCESS" if success else "FAILED"
                        args_text = json.dumps(call_info.get("args") or {}, ensure_ascii=False)
                        message_lines = [
                            f"[TOOL] {tool_name} execution {status_text}.",
                            f"[TOOL] Args: {args_text}",
                        ]
                        if last.get("output"):
                            message_lines.append(f"[TOOL] Output: {_truncate_text(last['output'], 120)}")
                        if last.get("error"):
                            message_lines.append(f"[TOOL] Error: {last['error']}")
                        message = "\n".join(message_lines) + "\n"
                        write_sentence(message, force_first=is_first_sentence)
                        is_first_sentence = False
                        full_response_text += message
                        append_to_buffer('assistant', message.strip())
                elif step == "__end__":
                    break
        except Exception as exc:
            util.log(1, f"执行工具工作流时出错: {exc}")
            if is_agent_think_start:
                closing = "</think>"
                accumulated_text += closing
                full_response_text += closing
            return False

        if final_state is None:
            if is_agent_think_start:
                closing = "</think>"
                accumulated_text += closing
                full_response_text += closing
            return False

        if not final_stream_done and is_agent_think_start:
            closing = "</think>"
            accumulated_text += closing
            full_response_text += closing
            util.log(1, f"工具工作流未能完成，状态: {final_state.get('status')}")

        final_state_messages = final_state.get("messages") if final_state else None
        if final_state_messages and len(final_state_messages) > len(messages_buffer):
            messages_buffer.extend(final_state_messages[len(messages_buffer):])
            if len(messages_buffer) > 60:
                del messages_buffer[:-60]

        return final_stream_done

    def run_direct_llm() -> bool:
        nonlocal full_response_text, accumulated_text, is_first_sentence, messages_buffer
        try:
            # 统一使用 _build_final_messages 构建消息，确保历史对话始终被包含
            summary_state: AgentState = {
                "request": content,
                "messages": messages_buffer,
                "tool_results": [],
                "planner_preview": None,
                "context": {
                    "system_prompt": system_prompt,
                    "knowledge_context": knowledge_context,
                    "observation": observation,
                },
            }

            final_messages = _build_final_messages(summary_state)
            stream_response_chunks(llm.stream(final_messages))
            return True
        except requests.exceptions.RequestException as exc:
            util.log(1, f"请求失败: {exc}")
            error_message = "抱歉，我现在太忙了，休息一会，请稍后再试。"
            write_sentence(error_message, force_first=is_first_sentence)
            is_first_sentence = False
            full_response_text = error_message
            accumulated_text = ""
            return False

    workflow_success = False
    if tool_registry:
        workflow_success = run_workflow(tool_registry)

    if (not tool_registry or not workflow_success) and not sm.should_stop_generation(username, conversation_id=conversation_id):
        run_direct_llm()

    if not sm.should_stop_generation(username, conversation_id=conversation_id):
        finalize_stream(force_end=True)

    if state_mgr is not None:
        try:
            state_mgr.end_session(username, conversation_id=conversation_id)
        except Exception:
            pass
    else:
        try:
            from utils.stream_state_manager import get_state_manager

            get_state_manager().end_session(username, conversation_id=conversation_id)
        except Exception:
            pass

    final_text = full_response_text.split("</think>")[-1] if full_response_text else ""

    # 使用新记忆系统异步处理agent回复
    try:
        import asyncio

        # 创建新的事件循环（在独立线程中运行）
        def async_memory_task():
            """在独立线程中运行异步记忆存储"""
            try:
                # 创建新的事件循环
                loop = asyncio.new_event_loop()
                asyncio.set_event_loop(loop)

                # 运行异步任务
                loop.run_until_complete(
                    memory_system.process_agent_reply_async(
                        final_text,
                        user_id=username,
                        current_user_content=content
                    )
                )

                # 关闭循环
                loop.close()
            except Exception as e:
                util.log(1, f"异步记忆存储失败: {e}")

        # 启动独立线程执行异步任务
        MyThread(target=async_memory_task).start()
        util.log(1, f"异步记忆存储任务已启动")

    except Exception as exc:
        util.log(1, f"异步记忆处理启动失败: {exc}")

    return final_text
def clear_agent_memory(username=None):
    """
    清除指定用户的记忆（使用新记忆系统）

    Args:
        username: 用户名，如果为None则清除当前用户的记忆

    Returns:
        bool: 是否清除成功
    """
    global memory_system, current_username

    try:
        # 确定要清除的用户ID
        user_id = username if username else current_username
        if not user_id:
            user_id = "User"  # 默认用户

        util.log(1, f"正在清除用户 {user_id} 的记忆...")

        # 调用新记忆系统的清除方法
        result = memory_system.clear_user_history(user_id=user_id)

        util.log(1, f"用户 {user_id} 的记忆清除完成: {result}")
        return True

    except Exception as e:
        util.log(1, f"清除用户记忆时出错: {str(e)}")
        return False

def get_mcp_tools() -> List[Dict[str, Any]]:
    """
    从共享缓存获取所有可用且已启用的MCP工具列表。
    """
    try:
        tools = mcp_tool_registry.get_enabled_tools()
        return tools or []
    except Exception as e:
        util.log(1, f"获取工具列表出错：{e}")
        return []


if __name__ == "__main__":
    # 记忆系统已在模块加载时初始化，无需再次调用
    for _ in range(3):
        query = "Who is Fay?"
        response = question(query, "User")
        print(f"Q: {query}")
        print(f"A: {response}")
        time.sleep(1)
